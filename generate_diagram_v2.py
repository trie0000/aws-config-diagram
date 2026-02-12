"""
generate_diagram_v2.py: AWS Config JSON -> AWS-style Network Diagram (pptx)

Generates a diagram matching AWS Architecture diagram style with:
- AWS official icons (PNG) for 30+ service types
- Gateway Column layout: IGW/ALB in dedicated VPC-level column
- Horizontal AZ rows (AZ-A top, AZ-C bottom)
- Dynamic subnet column widths based on icon count
- Bottom-up height calculation (prevents overflow)
- TCP(port) arrow labels on connections
- External actors (Internet, End User) and edge services (Route53, CloudFront, API GW)
- VPC-external services row (ACM, WAF, KMS, CloudTrail, CloudWatch)
- Serverless services row (Lambda, DynamoDB, SQS, SNS, S3)
- Same-AZ arrows preferred; cross-AZ arrows minimized

Usage:
    python generate_diagram_v2.py tabelog_aws_config.json

Version: 4.0.0
Last Updated: 2026-02-11
"""

import json
import sys
import os
from collections import defaultdict

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

from generate_diagram import AWSConfigParser

# ============================================================
# Icon paths
# ============================================================
ICON_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "icons")

ICONS = {}
_ICON_NAMES = [
    # Existing
    "internet", "users", "client", "ec2", "rds", "alb", "igw", "nat",
    "waf", "acm", "s3", "s3_bucket", "vpc_endpoint", "vpc_icon",
    "region", "aws_cloud", "public_subnet", "private_subnet",
    # New
    "lambda", "ecs", "eks", "autoscaling", "cloudfront", "apigateway",
    "route53", "dynamodb", "elasticache", "redshift", "sqs", "sns",
    "kms", "cloudtrail", "cloudwatch", "elasticbeanstalk",
]
for name in _ICON_NAMES:
    path = os.path.join(ICON_DIR, f"{name}.png")
    if os.path.exists(path):
        ICONS[name] = path


# ============================================================
# Colors
# ============================================================
class C:
    CLOUD_BG = RGBColor(0xF5, 0xF5, 0xF5)
    CLOUD_BD = RGBColor(0x9A, 0x9A, 0x9A)
    VPC_BG   = RGBColor(0xF0, 0xF7, 0xF0)
    VPC_BD   = RGBColor(0x1A, 0x8C, 0x1A)

    PUB_BG   = RGBColor(0xE8, 0xF5, 0xE2)
    PUB_BD   = RGBColor(0x24, 0x8F, 0x24)
    PRIV_BG  = RGBColor(0xE3, 0xED, 0xF7)
    PRIV_BD  = RGBColor(0x14, 0x7E, 0xBA)
    GW_BG    = RGBColor(0xFD, 0xF0, 0xE0)
    GW_BD    = RGBColor(0xCC, 0xBB, 0x99)

    ARROW_INET = RGBColor(0x00, 0x73, 0xBB)
    ARROW_AWS  = RGBColor(0xED, 0x7D, 0x1C)
    ARROW_GRAY = RGBColor(0x88, 0x88, 0x88)
    ARROW_PEER = RGBColor(0x00, 0x96, 0x88)

    TEXT   = RGBColor(0x33, 0x33, 0x33)
    TEXT_G = RGBColor(0x77, 0x77, 0x77)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


# ============================================================
# V2 Generator (v4.0 - Gateway Column layout)
# ============================================================
class DiagramV2:
    def __init__(self, parser: AWSConfigParser):
        self.p = parser
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)
        self.pos = {}       # key -> (cx, cy, hw, hh) bounding box
        self.shapes = {}    # key -> picture Shape (for connector binding)

    def _score_vpc(self, v):
        """Score a VPC by resource count (higher = more resources)."""
        vid = v["id"]
        score = (len(self.p.get_subnets_for_vpc(vid))
                 + len(self.p.get_albs_for_vpc(vid)) * 10
                 + len(self.p.get_rds_for_vpc(vid)) * 5)
        sub_ids = {s["id"] for s in self.p.get_subnets_for_vpc(vid)}
        for item in self.p.by_type["AWS::EC2::Instance"]:
            cfg = item.get("configuration", {})
            if cfg.get("subnetId", "") in sub_ids:
                score += 1
        return score

    def list_vpcs(self):
        """List all VPCs in the config with scores. Returns list of dicts."""
        vpcs = self.p.get_vpcs()
        result = []
        for v in vpcs:
            score = self._score_vpc(v)
            result.append({**v, "score": score})
        return result

    def generate(self, out, vpc_ids=None):
        """Generate diagram.

        Args:
            out: output pptx file path
            vpc_ids: list of VPC IDs to draw (each on its own slide).
                     If None, auto-selects the best VPC.
        """
        vpcs = self.p.get_vpcs()
        if not vpcs:
            print("No VPCs found")
            return

        if vpc_ids:
            # Draw specified VPCs
            vpc_map = {v["id"]: v for v in vpcs}
            targets = []
            for vid in vpc_ids:
                if vid in vpc_map:
                    targets.append(vpc_map[vid])
                else:
                    print(f"  Warning: VPC {vid} not found in config, skipping")
            if not targets:
                print("Error: none of the specified VPCs found")
                return
        else:
            # Auto-select: skip default, pick highest-scoring
            non_default = [v for v in vpcs if not v.get("is_default", False)]
            candidates = non_default if non_default else vpcs

            if len(candidates) == 1:
                targets = [candidates[0]]
            else:
                scored = [(self._score_vpc(v), v) for v in candidates]
                scored.sort(key=lambda x: -x[0])
                for score, v in scored:
                    print(f"  VPC {v['name']} ({v['id']}): score={score}")
                targets = [scored[0][1]]

        for v in targets:
            print(f"  Drawing VPC: {v['name']} ({v['id']})")
            self.pos = {}
            self.shapes = {}
            self._build(v)

        self.prs.save(out)
        print(f"Saved: {out}")

    # ==========================================================
    # Layout calculation (bottom-up to prevent overflow)
    # ==========================================================
    def _calc_layout(self, has_edge, has_peering, has_svless, has_infra,
                     tiers, azs, res_ctx, gw_item_count=0,
                     icon_conns=None, all_icon_tiers=None):
        """Calculate all layout positions. Returns dict of Inches values.

        Automatically expands VPC/AZ/Cloud boxes so that icons are never
        clipped.  The algorithm works bottom-up:

        1. Determine the minimum AZ row height from icon content.
        2. Determine the minimum GW column height from its item count.
        3. Take the larger of the two as the minimum AZ height.
        4. Compute VPC / Cloud heights, expanding if necessary (up to
           slide bottom limit).
        """
        L = {}

        # Left panel width
        L['left_w'] = Inches(2.3) if has_edge else Inches(1.3)
        L['right_margin'] = Inches(0.9) if has_peering else Inches(0.3)

        # AWS Cloud box
        L['cloud_x'] = L['left_w']
        L['cloud_y'] = Inches(0.35)
        L['cloud_w'] = Inches(16) - L['left_w'] - L['right_margin']
        L['cloud_bottom'] = Inches(8.95)

        # Bottom service rows (inside cloud, below VPC)
        bottom_row_h = Inches(0.65)
        n_bottom = (1 if has_infra else 0) + (1 if has_svless else 0)
        bottom_total = Inches(0.10) + bottom_row_h * n_bottom if n_bottom > 0 \
            else Inches(0.05)

        # --- Minimum AZ row height (initial estimate, 1 icon row) ---
        # sub_y offset(0.22) + header(0.22) + icon_row(0.75) + pad(0.05)
        az_gap = Inches(0.15)
        min_az_h_1row = Inches(1.05)

        # GW column minimum from item count
        gw_min_h = Inches(0.15) + Inches(0.78) * max(gw_item_count, 1)
        min_az_from_gw = (gw_min_h - az_gap) / 2
        min_az_h = max(min_az_h_1row, min_az_from_gw)

        # VPC box (inside Cloud)
        cloud_pad = Inches(0.15)
        L['vpc_x'] = L['cloud_x'] + cloud_pad
        L['vpc_y'] = L['cloud_y'] + Inches(0.40)
        L['vpc_w'] = L['cloud_w'] - 2 * cloud_pad

        vpc_header = Inches(0.30)
        vpc_pad = Inches(0.10)

        # --- First pass: compute column widths and max icon rows ---
        # Use initial AZ height to determine VPC/GW/subnet geometry
        available_vpc_h = (L['cloud_bottom'] - L['vpc_y']
                           - bottom_total - Inches(0.10))
        min_vpc_h = vpc_header + vpc_pad + 2 * min_az_h + az_gap
        L['vpc_h'] = max(available_vpc_h, min_vpc_h)

        L['gw_x'] = L['vpc_x'] + vpc_pad
        L['gw_w'] = Inches(1.30)

        col_gap = Inches(0.10)
        subnet_area_x = L['gw_x'] + L['gw_w'] + col_gap
        subnet_area_w = (L['vpc_x'] + L['vpc_w'] - vpc_pad) - subnet_area_x

        if icon_conns is None:
            icon_conns = {}
        if all_icon_tiers is None:
            all_icon_tiers = {}
        col_widths, max_icon_rows = self._calc_col_widths(
            tiers, azs, subnet_area_w, col_gap, res_ctx,
            icon_conns, all_icon_tiers)

        # --- Recalculate min AZ height from actual icon rows ---
        # sub_y offset(0.22) + header(0.22) + rows * 0.75 + pad(0.05)
        icon_row_h = Inches(0.75)
        content_min_az = Inches(0.22 + 0.22 + 0.05) + icon_row_h * max(max_icon_rows, 1)

        # Hard limit: AZ height must fit within slide bounds
        # slide_bottom(8.95) - cloud_y(0.35) - cloud_header(0.40)
        #   - vpc_header - vpc_pad - az_gap - bottom_total - margin(0.10)
        # = available for 2 AZ rows → max_az_from_slide = available / 2
        slide_bottom = Inches(8.95)
        max_az_from_slide = (slide_bottom - L['vpc_y'] - vpc_header
                             - vpc_pad - az_gap - bottom_total
                             - Inches(0.10)) / 2
        # If content needs more than slide allows, cap to slide limit
        content_min_az = min(content_min_az, max_az_from_slide)

        min_az_h = max(min_az_h, content_min_az)

        # Cap: content-driven max (avoid huge whitespace for 1-row case)
        max_az_h = content_min_az + Inches(0.20)  # small padding
        max_az_h = max(max_az_h, min_az_from_gw)  # but respect GW needs
        max_az_h = min(max_az_h, max_az_from_slide)  # never exceed slide

        # Recompute VPC height with updated min_az_h
        min_vpc_h = vpc_header + vpc_pad + 2 * min_az_h + az_gap
        L['vpc_h'] = max(available_vpc_h, min_vpc_h)

        # Ensure cloud_bottom never exceeds slide bottom
        actual_cloud_bottom = max(
            L['cloud_bottom'],
            L['vpc_y'] + L['vpc_h'] + bottom_total + Inches(0.10))
        actual_cloud_bottom = min(actual_cloud_bottom, slide_bottom)
        L['cloud_bottom'] = actual_cloud_bottom
        L['cloud_h'] = L['cloud_bottom'] - L['cloud_y']

        # Refit VPC height to cloud bounds
        max_vpc_h = L['cloud_bottom'] - L['vpc_y'] - bottom_total - Inches(0.10)
        L['vpc_h'] = min(L['vpc_h'], max_vpc_h)

        # AZ rows inside VPC
        az_area_top = L['vpc_y'] + vpc_header
        az_area_h = L['vpc_h'] - vpc_header - vpc_pad
        az_h = (az_area_h - az_gap) / 2
        az_h = min(az_h, max_az_h)  # cap to content-driven max

        L['az_a_y'] = az_area_top
        L['az_c_y'] = az_area_top + az_h + az_gap
        L['az_h'] = az_h

        # Gateway column
        L['gw_y'] = L['az_a_y']
        L['gw_h'] = (L['az_c_y'] + az_h) - L['az_a_y']

        # Subnet column positions
        L['pub_x'] = subnet_area_x
        L['pub_w'] = col_widths['Public']
        L['priv_x'] = L['pub_x'] + L['pub_w'] + col_gap
        L['priv_w'] = col_widths['Private']
        L['iso_x'] = L['priv_x'] + L['priv_w'] + col_gap
        L['iso_w'] = col_widths['Isolated']

        # Bottom service rows (inside cloud, below VPC)
        L['infra_y'] = L['vpc_y'] + L['vpc_h'] + Inches(0.15) \
            if has_infra else None
        if has_svless:
            if has_infra:
                L['svless_y'] = L['infra_y'] + bottom_row_h
            else:
                L['svless_y'] = L['vpc_y'] + L['vpc_h'] + Inches(0.15)
        else:
            L['svless_y'] = None

        return L

    def _calc_col_widths(self, tiers, azs, subnet_area_w, col_gap, res_ctx,
                          icon_conns, all_icon_tiers):
        """Calculate column widths proportional to max icon count per tier.

        Returns (col_widths, max_icon_rows):
          col_widths: {tier: width_emu}
          max_icon_rows: int — maximum number of icon rows across all subnets
        """
        icon_slot = Inches(1.10)
        min_col = Inches(1.80)

        # Collect icons per tier/az for counting and row computation
        all_subnet_icons = {}  # (tier, az_idx) -> icons list
        max_icons = {"Public": 0, "Private": 0, "Isolated": 0}
        for tier in max_icons:
            for ai, az in enumerate(azs):
                subs = tiers.get(tier, {}).get(az, [])
                if subs:
                    icons, _aux = self._collect_subnet_icons(tier, ai,
                                                              subs[0], res_ctx)
                    all_subnet_icons[(tier, ai)] = icons
                    max_icons[tier] = max(max_icons[tier], len(icons))

        # Compute desired widths
        active_tiers = [t for t in ["Public", "Private", "Isolated"]
                        if max_icons[t] > 0 or any(tiers.get(t, {}).get(az, []) for az in azs)]
        n_gaps = max(len(active_tiers) - 1, 0)
        available = subnet_area_w - col_gap * n_gaps

        desired = {}
        for t in ["Public", "Private", "Isolated"]:
            desired[t] = max(max_icons[t] * icon_slot, min_col)

        total_desired = sum(desired.values())
        if total_desired > 0 and available > 0:
            scale = available / total_desired
            col_widths = {t: desired[t] * scale for t, w in desired.items()}
        else:
            equal = available / 3 if available > 0 else min_col
            col_widths = {"Public": equal, "Private": equal, "Isolated": equal}

        # Calculate max icon rows using topology-based layout (exact count)
        max_icon_rows = 1
        for tier in ["Public", "Private", "Isolated"]:
            for ai in range(len(azs)):
                icons = all_subnet_icons.get((tier, ai), [])
                if icons:
                    rows = self._num_icon_rows(icons, col_widths[tier],
                                                tier, icon_conns,
                                                all_icon_tiers)
                    max_icon_rows = max(max_icon_rows, rows)

        return col_widths, max_icon_rows

    # ==========================================================
    # Collect icons for a subnet (used for counting & drawing)
    # ==========================================================
    def _collect_subnet_icons(self, tier, ai, sub, ctx):
        """Return (main_icons, aux_labels) for a subnet.

        main_icons: [(icon_name, label, key), ...] — direct data-path services
        aux_labels: [(icon_name, short_label, key), ...] — orchestration/management
                    services shown as small icons in subnet top-right corner
        """
        icons = []
        aux = []

        if tier == "Public":
            # NAT Gateway is placed on subnet border (not as regular icon)
            # EC2
            for inst in self.p.get_instances_for_subnet(sub["id"]):
                icons.append(("ec2", f"EC2\n{inst['name']}",
                              f"ec2_{inst['id']}"))

        elif tier == "Private":
            # EC2 — direct data-path
            for inst in self.p.get_instances_for_subnet(sub["id"]):
                icons.append(("ec2", f"EC2\n{inst['name']}",
                              f"ec2_{inst['id']}"))
            # Lambda (VPC-attached) — direct data-path
            for lf in ctx['lambdas_vpc']:
                if sub["id"] in lf.get("vpc_subnet_ids", []):
                    icons.append(("lambda", f"Lambda\n{lf['name'][:12]}",
                                  f"lambda_{lf['id']}"))
            # ECS — direct data-path (container service)
            for svc in ctx['ecs_services']:
                if sub["id"] in svc.get("subnet_ids", []):
                    icons.append(("ecs", f"ECS\n{svc['name']}",
                                  f"ecs_{svc['id']}"))
            # EKS — direct data-path (container service)
            for ek in ctx['eks_clusters']:
                if sub["id"] in ek.get("subnet_ids", []):
                    icons.append(("eks", f"EKS\n{ek['name'][:12]}",
                                  f"eks_{ek['id']}"))

            # --- Management/orchestration → aux badges ---
            # ElasticBeanstalk (AZ-A only)
            if ai == 0:
                for eb in ctx['eb_envs']:
                    aux.append(("elasticbeanstalk",
                               f"Beanstalk\n{eb['name'][:12]}",
                               f"eb_{eb['id']}"))

        elif tier == "Isolated":
            # RDS
            for db in ctx['rdss']:
                if sub["id"] in db["subnet_ids"]:
                    role = "(Primary)" if ai == 0 else "(Standby)"
                    icons.append(("rds", f"Amazon RDS\n{role}",
                                  f"rds_{db['id']}_{ai}"))
            # ElastiCache
            for cc in ctx['cache_clusters']:
                icons.append(("elasticache", f"ElastiCache\n{cc['engine']}",
                              f"cache_{cc['id']}"))
            # Redshift
            for rc in ctx['rs_clusters']:
                icons.append(("redshift", f"Redshift\n{rc['name'][:10]}",
                              f"redshift_{rc['id']}"))

        return icons, aux

    # ==========================================================
    # Connection-aware icon placement
    # ==========================================================
    def _build_icon_connections(self, sg_conns, svc_conns, sg_map):
        """Pre-compute bidirectional connection graph for icon keys.

        Returns dict: icon_key -> set of connected icon_keys.
        Keys use the same naming as _collect_subnet_icons (ec2_xxx, rds_xxx_0, etc).
        """
        conns = defaultdict(set)

        # SG-based connections
        for conn in sg_conns:
            frs = sg_map.get(conn["from_sg"], [])
            trs = sg_map.get(conn["to_sg"], [])
            for fr in frs:
                fk = f"{fr['prefix']}{fr['id']}"
                for tr in trs:
                    tk = f"{tr['prefix']}{tr['id']}"
                    if tr['type'] == 'RDS':
                        # RDS uses _0/_1 suffixes per AZ
                        conns[fk].add(tk + "_0")
                        conns[fk].add(tk + "_1")
                        conns[tk + "_0"].add(fk)
                        conns[tk + "_1"].add(fk)
                    else:
                        conns[fk].add(tk)
                        conns[tk].add(fk)

        # ALB connections (ALB key format: alb_xxx)
        for conn in sg_conns:
            frs = sg_map.get(conn["from_sg"], [])
            trs = sg_map.get(conn["to_sg"], [])
            for fr in frs:
                if fr['type'] == 'ALB':
                    fk = f"alb_{fr['id']}"
                    for tr in trs:
                        tk = f"{tr['prefix']}{tr['id']}"
                        conns[fk].add(tk)
                        conns[tk].add(fk)

        # Service connections
        for conn in svc_conns:
            ft = conn.get("from_type", "")
            fi = conn.get("from_id", "")
            tt = conn.get("to_type", "")
            ti = conn.get("to_id", "")
            fk = ft if ft else f"{ft}_{fi}"
            tk = tt if tt else f"{tt}_{ti}"
            conns[fk].add(tk)
            conns[tk].add(fk)

        return conns

    # Tier ordering for left-to-right flow scoring
    _TIER_ORDER = {"Gateway": 0, "Public": 1, "Private": 2, "Isolated": 3,
                   "external": 4}

    def _sort_icons_by_connections(self, icons, tier, icon_conns, all_icon_tiers):
        """Sort icons within a subnet so connected-to-left appear left, etc.

        Icons that receive connections from a tier to the LEFT should be placed
        leftward.  Icons that send connections to a tier to the RIGHT should be
        placed rightward.  This minimises arrow length across the diagram.
        """
        if len(icons) <= 1:
            return icons

        my_order = self._TIER_ORDER.get(tier, 2)

        def score(icon_tuple):
            """Lower score = place more to the left."""
            _icon, _label, key = icon_tuple
            neighbors = icon_conns.get(key, set())
            if not neighbors:
                return 0  # no connections → neutral, keep original order

            # Average tier order of connected services
            tier_scores = []
            for nk in neighbors:
                n_tier = all_icon_tiers.get(nk)
                if n_tier is not None:
                    tier_scores.append(self._TIER_ORDER.get(n_tier, 2))
            if not tier_scores:
                return 0
            avg = sum(tier_scores) / len(tier_scores)
            # Icons connected to LEFT tiers (lower order) get lower score (placed left)
            # Icons connected to RIGHT tiers (higher order) get higher score (placed right)
            return avg - my_order

        return sorted(icons, key=score)

    @staticmethod
    def _icons_per_row(n_icons, area_w):
        """How many icons fit in one row at ideal spacing (1.10in).

        Uses ideal spacing as the cap so icons aren't crammed into one row.
        This ensures multi-row layout when there are many icons.
        """
        margin = Inches(0.10)
        available = area_w - 2 * margin
        icon_slot_ideal = Inches(1.10)
        if n_icons == 0:
            return 1
        # Max icons that fit at ideal spacing
        per_row = max(1, int(available / icon_slot_ideal))
        return min(per_row, n_icons)

    def _build_topology_rows(self, icons, area_w, tier, icon_conns,
                             all_icon_tiers):
        """Build row layout using topology-based BFS chaining.

        Returns list of rows, where each row is a list of (icon, label, key).
        Used by both _num_icon_rows (for height calculation) and
        _place_icons_grid (for actual drawing).
        """
        if not icons:
            return []

        margin = Inches(0.10)
        spacing = Inches(1.10)
        available = area_w - 2 * margin
        max_per_row = max(1, int(available / spacing))

        my_order = self._TIER_ORDER.get(tier, 2)
        icon_keys = [ic[2] for ic in icons]
        icon_map = {ic[2]: ic for ic in icons}
        icon_set = set(icon_keys)
        placed = set()
        rows = []

        # --- Entry points: receive from LEFT tiers ---
        def left_score(key):
            neighbors = icon_conns.get(key, set())
            scores = []
            for nk in neighbors:
                nt = all_icon_tiers.get(nk)
                if nt is not None:
                    scores.append(self._TIER_ORDER.get(nt, 2))
            return min(scores) if scores else my_order

        entry_keys = []
        other_keys = []
        for key in icon_keys:
            neighbors = icon_conns.get(key, set())
            has_left = any(
                self._TIER_ORDER.get(all_icon_tiers.get(nk, ""), 99) < my_order
                for nk in neighbors
                if all_icon_tiers.get(nk) is not None
            )
            if has_left:
                entry_keys.append(key)
            else:
                other_keys.append(key)

        entry_keys.sort(key=left_score)

        # --- BFS chain builder ---
        def build_chain(start_key):
            chain = [start_key]
            placed.add(start_key)
            frontier = [start_key]
            while frontier:
                cur = frontier.pop(0)
                neighbors = icon_conns.get(cur, set())
                for nk in sorted(neighbors):
                    if nk in icon_set and nk not in placed:
                        chain.append(nk)
                        placed.add(nk)
                        frontier.append(nk)
            return chain

        # Build chains from entry points
        multi_chains = []   # chains with BFS connections (>1 icon)
        single_icons = []   # isolated icons (no BFS neighbours in this subnet)

        for ek in entry_keys:
            if ek in placed:
                continue
            chain = build_chain(ek)
            if len(chain) > 1:
                multi_chains.append(chain)
            else:
                single_icons.append(chain[0])

        # Remaining unplaced icons
        for ok in other_keys:
            if ok in placed:
                continue
            chain = build_chain(ok)
            if len(chain) > 1:
                multi_chains.append(chain)
            else:
                single_icons.append(chain[0])

        # Emit multi-icon chains first (preserve BFS order within chain)
        for chain in multi_chains:
            for i in range(0, len(chain), max_per_row):
                row = [icon_map[k] for k in chain[i:i+max_per_row]]
                rows.append(row)

        # Pack isolated icons together into rows (fill each row to max)
        if single_icons:
            for i in range(0, len(single_icons), max_per_row):
                row = [icon_map[k] for k in single_icons[i:i+max_per_row]]
                rows.append(row)

        return rows

    def _num_icon_rows(self, icons, area_w, tier, icon_conns, all_icon_tiers):
        """Exact number of rows for topology-based layout."""
        if not icons:
            return 0
        rows = self._build_topology_rows(icons, area_w, tier,
                                          icon_conns, all_icon_tiers)
        return len(rows)

    def _place_icons_grid(self, sl, icons, area_x, area_w, y_base,
                          tier, icon_conns, all_icon_tiers):
        """Place icons using topology-based layout.

        Rules (from user):
        - If arrow doesn't cross another icon → place right
        - If right is full → place below
        - If right and below are full → diagonal right-down
        """
        rows = self._build_topology_rows(icons, area_w, tier,
                                          icon_conns, all_icon_tiers)
        if not rows:
            return

        margin = Inches(0.10)
        spacing = Inches(1.10)
        available = area_w - 2 * margin
        row_h = Inches(0.75)

        for row_i, row_icons in enumerate(rows):
            nr = len(row_icons)
            if nr * spacing <= available:
                sp = spacing
            else:
                sp = available / nr
            total = nr * sp
            start_x = area_x + margin + (available - total) / 2
            row_y = y_base + row_i * row_h

            for i, (icon_name, label, key) in enumerate(row_icons):
                x = int(start_x + i * sp)
                self._ibox(sl, x, int(row_y), icon_name, label, key)

    def _place_aux_badges(self, sl, aux, area_x, area_w, sub_y):
        """Place auxiliary service icons (small) in subnet top-right."""
        if not aux:
            return
        isz = Inches(0.28)           # small icon
        lbl_w = Inches(0.80)         # label width
        lbl_h = Inches(0.20)         # label height
        item_h = isz + lbl_h         # total per item
        gap = Inches(0.02)
        # right-align in subnet area
        total_w = lbl_w
        x = int(area_x + area_w - total_w - Inches(0.05))
        y = int(sub_y + Inches(0.03))
        for icon_name, label, key in aux:
            # icon centred above label
            ix = int(x + total_w / 2 - isz / 2)
            aux_pic = None
            if icon_name in ICONS:
                aux_pic = sl.shapes.add_picture(ICONS[icon_name], ix, y,
                                                isz, isz)
            # label below icon
            self._txt(sl, x, int(y + isz + Inches(0.01)),
                      lbl_w, lbl_h, label, 5, True, C.TEXT, PP_ALIGN.CENTER)
            # register pos for potential arrows
            icon_cx = int(x + total_w / 2)
            icon_cy = int(y + isz / 2)
            icon_hw = int(isz / 2)
            icon_hh = int(isz / 2)
            self.pos[key] = (icon_cx, icon_cy, icon_hw, icon_hh)
            if aux_pic is not None:
                self.shapes[key] = aux_pic
            y += int(item_h + gap)

    # ==========================================================
    # Main build
    # ==========================================================
    def _build(self, vpc):
        sl = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # ---- Gather all resources ----
        subs = self.p.get_subnets_for_vpc(vpc["id"])
        igw = self.p.get_igw_for_vpc(vpc["id"])
        nats = self.p.get_nat_gateways_for_vpc(vpc["id"])
        albs = self.p.get_albs_for_vpc(vpc["id"])
        rdss = self.p.get_rds_for_vpc(vpc["id"])
        s3s = self.p.get_s3_buckets()
        sg_conns = self.p.get_sg_connections()

        lambdas = self.p.get_lambda_functions()
        ecs_services = self.p.get_ecs_services()
        eks_clusters = self.p.get_eks_clusters()
        cf_dists = self.p.get_cloudfront_distributions()
        api_gws = self.p.get_api_gateways()
        r53_zones = self.p.get_route53_hosted_zones()
        dynamo_tables = self.p.get_dynamodb_tables()
        cache_clusters = self.p.get_elasticache_clusters()
        rs_clusters = self.p.get_redshift_clusters()
        sqs_queues = self.p.get_sqs_queues()
        sns_topics = self.p.get_sns_topics()
        kms_keys = self.p.get_kms_keys()
        ct_trails = self.p.get_cloudtrail_trails()
        cw_alarms = self.p.get_cloudwatch_alarms()
        svc_conns = self.p.get_service_connections()
        asgs = self.p.get_autoscaling_groups()
        eb_envs = self.p.get_elasticbeanstalk_environments()
        peerings = self.p.get_peering_connections()

        waf = None
        for a in albs:
            waf = self.p.get_waf_for_alb(a["id"])
            if waf:
                break

        lambdas_vpc = [l for l in lambdas if l.get("in_vpc")]
        lambdas_serverless = [l for l in lambdas if not l.get("in_vpc")]

        asg_by_subnet = {}
        for asg in asgs:
            for sid in asg.get("subnet_ids", []):
                asg_by_subnet[sid] = asg

        azs = sorted(set(s["az"] for s in subs if s["az"]))[:2]
        tiers = defaultdict(lambda: defaultdict(list))
        for s in subs:
            tiers[s["tier"]][s["az"]].append(s)

        # ---- Detect what services exist ----
        has_edge = bool(r53_zones or cf_dists or api_gws)
        has_peering = bool(peerings)

        # ALB-related services → place in Gateway Column (near ALB)
        # Order: WAF before ALB (filters traffic), ACM after ALB (provides TLS)
        gw_before_alb = []   # placed above ALB in gateway column
        if waf:
            gw_before_alb.append(("waf", "AWS WAF", "waf"))
        gw_after_alb = []    # placed below ALB in gateway column

        # VPC-external services → place in bottom row
        infra_items = []
        if kms_keys:
            infra_items.append(("kms", "AWS KMS", "kms"))
        if ct_trails:
            infra_items.append(("cloudtrail", "CloudTrail", "cloudtrail"))
        if cw_alarms:
            infra_items.append(("cloudwatch", "CloudWatch", "cloudwatch"))

        svless_items = []
        if lambdas_serverless:
            svless_items.append(("lambda", "Lambda", "lambda_svless"))
        if dynamo_tables:
            svless_items.append(("dynamodb", "DynamoDB", "dynamodb"))
        if sqs_queues:
            svless_items.append(("sqs", "Amazon SQS", "sqs"))
        if sns_topics:
            svless_items.append(("sns", "Amazon SNS", "sns"))
        if s3s:
            svless_items.append(("s3_bucket", "Amazon S3", "s3"))

        has_svless = bool(svless_items)
        has_infra = bool(infra_items)

        # Resource context for icon collection
        res_ctx = {
            'nats': nats, 'asg_by_subnet': asg_by_subnet,
            'eb_envs': eb_envs, 'ecs_services': ecs_services,
            'lambdas_vpc': lambdas_vpc, 'eks_clusters': eks_clusters,
            'rdss': rdss, 'cache_clusters': cache_clusters,
            'rs_clusters': rs_clusters,
        }

        # ===== PRE-COMPUTE CONNECTION GRAPH (before layout calc) =====
        sg_map = self.p.build_sg_to_resources_map()
        icon_conns = self._build_icon_connections(sg_conns, svc_conns, sg_map)

        # Build icon_key -> tier map for connection-aware placement
        all_icon_tiers = {}
        # Gateway Column services
        for alb in albs:
            all_icon_tiers[f"alb_{alb['id']}"] = "Gateway"
        if waf:
            all_icon_tiers["waf"] = "Gateway"
        # Subnet services (scan all AZs)
        for ai, az in enumerate(azs):
            for tier_name in ["Public", "Private", "Isolated"]:
                sub_list = tiers.get(tier_name, {}).get(az, [])
                if sub_list:
                    icons, _aux = self._collect_subnet_icons(
                        tier_name, ai, sub_list[0], res_ctx)
                    for _icon, _label, key in icons:
                        all_icon_tiers[key] = tier_name
        # External/serverless
        for _icon, _label, key in infra_items + svless_items:
            all_icon_tiers[key] = "external"

        # ===== CALCULATE LAYOUT (uses connection graph for row estimation) =====
        gw_item_count = len(gw_before_alb) + len(albs) + len(gw_after_alb)
        L = self._calc_layout(has_edge, has_peering, has_svless, has_infra,
                              tiers, azs, res_ctx, gw_item_count,
                              icon_conns, all_icon_tiers)

        # ===== DRAW STRUCTURE =====

        # Title
        self._txt(sl, Inches(0.2), Inches(0.05), Inches(6), Inches(0.3),
                  "AWS Network Architecture", 14, True)

        # AWS Cloud box
        self._box(sl, L['cloud_x'], L['cloud_y'], L['cloud_w'], L['cloud_h'],
                  C.CLOUD_BG, C.CLOUD_BD)
        self._ilabel(sl, L['cloud_x'] + Inches(0.08),
                     L['cloud_y'] + Inches(0.05),
                     "aws_cloud", "AWS Cloud", 9, True)
        self._ilabel(sl, L['cloud_x'] + Inches(1.8),
                     L['cloud_y'] + Inches(0.05),
                     "region", f"Region: {vpc['region']}", 8, color=C.TEXT_G)

        # VPC box
        self._box(sl, L['vpc_x'], L['vpc_y'], L['vpc_w'], L['vpc_h'],
                  C.VPC_BG, C.VPC_BD, Pt(2))
        self._ilabel(sl, L['vpc_x'] + Inches(0.08),
                     L['vpc_y'] + Inches(0.05),
                     "vpc_icon", f"VPC  {vpc['cidr']}", 9, True,
                     color=C.VPC_BD)

        # Gateway Column (spans both AZ rows)
        self._box(sl, L['gw_x'], L['gw_y'], L['gw_w'], L['gw_h'],
                  C.GW_BG, C.GW_BD, Pt(0.75), 0.02)

        # Place ALB and related services in gateway column
        gw_icon_x = int(L['gw_x'] + L['gw_w'] / 2 - Inches(0.6))
        gw_cursor_y = L['gw_y'] + Inches(0.15)

        # Track Y positions for arrow alignment
        gw_first_y = int(gw_cursor_y)  # Y of first service in GW column
        alb_y_pos = None

        # WAF before ALB (filters incoming traffic)
        for icon, label, key in gw_before_alb:
            self._ibox(sl, gw_icon_x, int(gw_cursor_y),
                       icon, label, key, nobg=True)
            gw_cursor_y += Inches(0.85)

        for alb in albs:
            alb_y_pos = int(gw_cursor_y)
            self._ibox(sl, gw_icon_x, int(gw_cursor_y),
                       "alb", "Elastic Load\nBalancing",
                       f"alb_{alb['id']}")
            gw_cursor_y += Inches(1.0)

        # ACM after ALB (provides TLS cert)
        for icon, label, key in gw_after_alb:
            self._ibox(sl, gw_icon_x, int(gw_cursor_y),
                       icon, label, key, nobg=True)
            gw_cursor_y += Inches(0.85)

        # If no ALB, use first GW service Y
        if alb_y_pos is None:
            alb_y_pos = gw_first_y

        # ===== IGW straddling VPC left border =====
        # Align Y with first GW column service so arrow is horizontal
        if igw:
            igw_isz = Inches(0.42)
            igw_tw = Inches(1.2)
            igw_x = int(L['vpc_x'] - igw_tw / 2)
            igw_y = gw_first_y  # aligned with WAF or ALB
            igw_ix = int(igw_x + igw_tw / 2 - igw_isz / 2)
            igw_pic = None
            if "igw" in ICONS:
                igw_pic = sl.shapes.add_picture(ICONS["igw"], igw_ix, igw_y,
                                                igw_isz, igw_isz)
            self._txt(sl, igw_x, int(igw_y + igw_isz + Inches(0.01)),
                      igw_tw, Inches(0.28), "Internet\nGateway",
                      6, True, C.TEXT, PP_ALIGN.CENTER)
            igw_key = f"igw_{igw['id']}"
            cx = int(igw_x + igw_tw / 2)
            cy = int(igw_y + igw_isz / 2)
            hw = int(igw_isz / 2)
            hh = int(igw_isz / 2)
            self.pos[igw_key] = (cx, cy, hw, hh)
            if igw_pic is not None:
                self.shapes[igw_key] = igw_pic

        # ===== AZ ROWS WITH SUBNET COLUMNS =====
        for ai, az in enumerate(azs):
            row_y = L['az_a_y'] if ai == 0 else L['az_c_y']
            az_h = L['az_h']
            az_short = az.split("-")[-1].upper() if "-" in az else az.upper()

            # AZ label (placed to the right of gateway column)
            self._txt(sl, L['pub_x'], row_y + Inches(0.02),
                      Inches(3.5), Inches(0.2),
                      f"Availability Zone {az_short}", 8, True, C.TEXT_G)

            sub_y = row_y + Inches(0.22)
            sub_h = az_h - Inches(0.27)

            # Icon Y: place just below subnet header
            icon_y_base = sub_y + Inches(0.22)

            # --- Public Subnet ---
            ps = tiers.get("Public", {}).get(az, [])
            if ps:
                sub = ps[0]
                self._box(sl, L['pub_x'], sub_y, L['pub_w'], sub_h,
                          C.PUB_BG, C.PUB_BD)
                self._ilabel(sl, L['pub_x'] + Inches(0.05),
                             sub_y + Inches(0.03),
                             "public_subnet",
                             f"Public subnet  {sub['cidr']}", 7,
                             color=C.PUB_BD)
                icons, aux = self._collect_subnet_icons("Public", ai, sub,
                                                        res_ctx)
                self._place_icons_grid(sl, icons, L['pub_x'], L['pub_w'],
                                       icon_y_base, "Public",
                                       icon_conns, all_icon_tiers)
                self._place_aux_badges(sl, aux, L['pub_x'], L['pub_w'],
                                       sub_y)
                # NAT Gateway straddling top border of Public Subnet
                if ai == 0:
                    for nat in nats:
                        if nat["subnet_id"] == sub["id"]:
                            nat_isz = Inches(0.42)
                            nat_tw = Inches(1.2)
                            # Centre icon on the top border line of subnet
                            nat_x = int(L['pub_x'] + L['pub_w']
                                        - nat_tw - Inches(0.05))
                            nat_y = int(sub_y - nat_isz / 2)
                            nat_ix = int(nat_x + nat_tw / 2
                                         - nat_isz / 2)
                            nat_pic = None
                            if "nat" in ICONS:
                                nat_pic = sl.shapes.add_picture(
                                    ICONS["nat"], nat_ix, nat_y,
                                    nat_isz, nat_isz)
                            self._txt(sl, nat_x,
                                      int(nat_y + nat_isz + Inches(0.01)),
                                      nat_tw, Inches(0.28),
                                      f"NAT GW\n{nat['public_ip']}",
                                      6, True, C.TEXT, PP_ALIGN.CENTER)
                            nk = f"nat_{nat['id']}"
                            cx = int(nat_x + nat_tw / 2)
                            cy = int(nat_y + nat_isz / 2)
                            hw = int(nat_isz / 2)
                            hh = int(nat_isz / 2)
                            self.pos[nk] = (cx, cy, hw, hh)
                            if nat_pic is not None:
                                self.shapes[nk] = nat_pic

            # --- Private Subnet ---
            pvs = tiers.get("Private", {}).get(az, [])
            if pvs:
                sub = pvs[0]
                self._box(sl, L['priv_x'], sub_y, L['priv_w'], sub_h,
                          C.PRIV_BG, C.PRIV_BD)
                self._ilabel(sl, L['priv_x'] + Inches(0.05),
                             sub_y + Inches(0.03),
                             "private_subnet",
                             f"Private subnet  {sub['cidr']}", 7,
                             color=C.PRIV_BD)
                icons, aux = self._collect_subnet_icons("Private", ai, sub,
                                                        res_ctx)
                self._place_icons_grid(sl, icons, L['priv_x'], L['priv_w'],
                                       icon_y_base, "Private",
                                       icon_conns, all_icon_tiers)
                self._place_aux_badges(sl, aux, L['priv_x'], L['priv_w'],
                                       sub_y)

            # --- Isolated Subnet ---
            isos = tiers.get("Isolated", {}).get(az, [])
            if isos:
                sub = isos[0]
                self._box(sl, L['iso_x'], sub_y, L['iso_w'], sub_h,
                          C.PRIV_BG, C.PRIV_BD)
                self._ilabel(sl, L['iso_x'] + Inches(0.05),
                             sub_y + Inches(0.03),
                             "private_subnet",
                             f"Private subnet  {sub['cidr']}", 7,
                             color=C.PRIV_BD)
                icons, aux = self._collect_subnet_icons("Isolated", ai, sub,
                                                        res_ctx)
                self._place_icons_grid(sl, icons, L['iso_x'], L['iso_w'],
                                       icon_y_base, "Isolated",
                                       icon_conns, all_icon_tiers)
                self._place_aux_badges(sl, aux, L['iso_x'], L['iso_w'],
                                       sub_y)

        # ===== External actors (left of Cloud) =====
        ext_x = int(Inches(0.3))

        # End User — align Y with IGW (= gw_first_y) for horizontal arrow
        user_y = gw_first_y
        self._ibox(sl, ext_x, user_y,
                   "users", "End User", "user", nobg=True)

        # Edge services — align Y with their arrow targets
        if has_edge:
            edge_gap = Inches(0.85)
            # CloudFront/Route53 → ALB: align with ALB Y
            # API GW → Lambda/ALB: align below
            edge_items = []
            if r53_zones:
                edge_items.append(("route53", "Route 53", "route53"))
            if cf_dists:
                edge_items.append(("cloudfront", "CloudFront", "cloudfront"))
            if api_gws:
                edge_items.append(("apigateway", "API Gateway", "apigateway"))

            if len(edge_items) == 1:
                # Single edge service: align with ALB
                icon, label, key = edge_items[0]
                self._ibox(sl, ext_x, alb_y_pos,
                           icon, label, key, nobg=True)
            else:
                # Multiple: centre the group around ALB Y
                total_h = edge_gap * (len(edge_items) - 1)
                start_y = int(alb_y_pos - total_h / 2)
                # Clamp to cloud top
                min_y = int(L['cloud_y'] + Inches(0.1))
                if start_y < min_y:
                    start_y = min_y
                for ei, (icon, label, key) in enumerate(edge_items):
                    self._ibox(sl, ext_x, int(start_y + edge_gap * ei),
                               icon, label, key, nobg=True)

        # ===== VPC-external services (below VPC, inside Cloud) =====
        # Left-aligned with fixed spacing (1.40in per icon slot)
        bottom_svc_gap = Inches(1.40)
        bottom_start_x = L['vpc_x'] + Inches(0.2)

        if infra_items and L['infra_y'] is not None:
            for idx, (icon, label, key) in enumerate(infra_items):
                self._ibox(sl,
                           int(bottom_start_x + bottom_svc_gap * idx),
                           int(L['infra_y']),
                           icon, label, key, nobg=True)

        if has_svless and L['svless_y'] is not None:
            for idx, (icon, label, key) in enumerate(svless_items):
                self._ibox(sl,
                           int(bottom_start_x + bottom_svc_gap * idx),
                           int(L['svless_y']),
                           icon, label, key, nobg=True)

        # ===== VPC Peering (straddling VPC right border) =====
        if peerings:
            # _ibox uses tw=1.2in with icon centred; place so icon
            # centre sits on the VPC right border line.
            ibox_tw = Inches(1.2)
            peer_x = int(L['vpc_x'] + L['vpc_w'] - ibox_tw / 2)
            # Clamp so the ibox doesn't exceed slide right edge
            slide_w = Inches(16)
            max_x = int(slide_w - ibox_tw - Inches(0.05))
            if peer_x > max_x:
                peer_x = max_x
            peer_y = int(L['vpc_y'] + L['vpc_h'] / 2 - Inches(0.24))
            for pi, peer in enumerate(peerings):
                if peer["requester_vpc"] == vpc["id"]:
                    peer_label = f"VPC Peering\n{peer['accepter_cidr']}"
                else:
                    peer_label = f"VPC Peering\n{peer['requester_cidr']}"
                self._ibox(sl, peer_x, int(peer_y + pi * Inches(1.0)),
                           "vpc_icon", peer_label,
                           f"peering_{peer['id']}", nobg=True)

        # ===== Legend =====
        legend_h = Inches(1.3)
        legend_y = max(user_y + int(Inches(1.0)), int(Inches(7.3)))
        # Clamp so legend + ASG annotation fit within slide bottom
        slide_bottom = int(Inches(8.95))
        asg_h = Inches(0.15) * (len(asgs) + 1) if asgs else 0
        total_needed = int(legend_h + Inches(0.10) + asg_h)
        if legend_y + total_needed > slide_bottom:
            legend_y = slide_bottom - total_needed
        self._legend(sl, Inches(0.15), legend_y)

        # ===== ASG Annotation (below legend, left panel) =====
        if asgs:
            asg_note_y = legend_y + int(legend_h) + int(Inches(0.10))
            asg_lines = ["Auto Scaling Groups:"]
            for asg in asgs:
                name = asg.get("name", asg["id"][:20])
                asg_lines.append(
                    f"  {name}  ({asg['min_size']}-{asg['max_size']})")
            self._txt(sl, Inches(0.15), asg_note_y,
                      Inches(2.2), Inches(0.15) * len(asg_lines),
                      "\n".join(asg_lines), 6, False, C.TEXT_G)

        # ===== Arrows =====
        inet_sgs = self.p.get_internet_facing_sgs()
        self._draw_arrows(sl, albs, nats, rdss, igw, sg_conns, s3s, waf, azs,
                          cf_dists, api_gws, r53_zones, lambdas_serverless,
                          svc_conns, inet_sgs)

    # ==========================================================
    # Arrow drawing - same-AZ preferred
    # ==========================================================
    def _draw_arrows(self, sl, albs, nats, rdss, igw, sg_conns, s3s, waf, azs,
                     cf_dists, api_gws, r53_zones, lambdas_svless, svc_conns,
                     inet_sgs=None):
        sg_map = self.p.build_sg_to_resources_map()
        drawn = set()

        # ---- Edge service chain ----
        if "route53" in self.pos and "cloudfront" in self.pos:
            self._arr(sl, "route53", "cloudfront", C.ARROW_INET, "DNS")
        elif "route53" in self.pos:
            for alb in albs:
                ak = f"alb_{alb['id']}"
                if ak in self.pos:
                    self._arr(sl, "route53", ak, C.ARROW_INET, "DNS")
                    break
            else:
                if igw:
                    igw_key = f"igw_{igw['id']}"
                    if igw_key in self.pos:
                        self._arr(sl, "route53", igw_key, C.ARROW_INET, "DNS")

        if "cloudfront" in self.pos:
            for alb in albs:
                ak = f"alb_{alb['id']}"
                if ak in self.pos:
                    self._arr(sl, "cloudfront", ak, C.ARROW_INET, "HTTPS")
                    break
            else:
                if igw:
                    igw_key = f"igw_{igw['id']}"
                    if igw_key in self.pos:
                        self._arr(sl, "cloudfront", igw_key, C.ARROW_INET,
                                  "HTTPS")

        if "apigateway" in self.pos:
            if "lambda_svless" in self.pos:
                self._arr(sl, "apigateway", "lambda_svless", C.ARROW_AWS,
                          "invoke")
            else:
                for alb in albs:
                    ak = f"alb_{alb['id']}"
                    if ak in self.pos:
                        self._arr(sl, "apigateway", ak, C.ARROW_AWS, "HTTP")
                        break

        # ---- End User -> IGW -> (WAF ->) ALB ----
        if igw:
            igw_key = f"igw_{igw['id']}"
            if igw_key in self.pos:
                self._arr(sl, "user", igw_key, C.ARROW_INET, "HTTPS")
                if waf and "waf" in self.pos:
                    # IGW -> WAF -> ALB (WAF is above ALB in gateway column)
                    self._arr(sl, igw_key, "waf", C.ARROW_INET, "TCP(80,443)")
                else:
                    for alb in albs:
                        alb_key = f"alb_{alb['id']}"
                        if alb_key in self.pos:
                            self._arr(sl, igw_key, alb_key, C.ARROW_INET,
                                      "TCP(80,443)")

        # ---- IGW -> internet-facing resources (0.0.0.0/0 inbound) ----
        if igw and inet_sgs:
            igw_key = f"igw_{igw['id']}"
            if igw_key in self.pos:
                # Collect SG IDs already covered (ALB SGs)
                alb_sg_ids = set()
                for alb in albs:
                    alb_sg_ids.update(alb.get("sg_ids", []))

                for isg in inet_sgs:
                    sg_id = isg["sg_id"]
                    if sg_id in alb_sg_ids:
                        continue  # ALB already has IGW->WAF->ALB path
                    port = isg.get("port", "")
                    proto = isg.get("protocol", "tcp").upper()
                    label = f"{proto}({port})" if port else ""
                    # Find resources in this SG
                    for res in sg_map.get(sg_id, []):
                        tk = self._resolve_key(
                            f"{res['prefix']}{res['id']}")
                        if tk and tk in self.pos:
                            aid = f"{igw_key}->{tk}"
                            if aid not in drawn:
                                drawn.add(aid)
                                self._arr(sl, igw_key, tk,
                                          C.ARROW_INET, label)

        # ---- SG-based internal connections ----
        for conn in sg_conns:
            frs = sg_map.get(conn["from_sg"], [])
            trs = sg_map.get(conn["to_sg"], [])
            port = conn.get("port", "")
            proto = conn.get("protocol", "tcp").upper()
            label = f"{proto}({port})" if port else ""

            fr_by_az = self._group_by_az(frs)
            tr_by_az = self._group_by_az(trs)

            # Same-AZ connections
            for az_suffix, az_frs in fr_by_az.items():
                if az_suffix == "_global":
                    continue
                az_trs = tr_by_az.get(az_suffix, [])
                for fr in az_frs:
                    for tr in az_trs:
                        fk = self._resolve_key(f"{fr['prefix']}{fr['id']}")
                        tk = self._resolve_key(f"{tr['prefix']}{tr['id']}")
                        if not fk or not tk:
                            continue
                        aid = f"{fk}->{tk}"
                        if aid not in drawn:
                            drawn.add(aid)
                            self._arr(sl, fk, tk, C.ARROW_AWS, label)

            # Global resources (ALB) -> all targets
            for fr in fr_by_az.get("_global", []):
                fk = self._resolve_key(f"{fr['prefix']}{fr['id']}")
                if not fk:
                    continue
                for tr in trs:
                    tk = self._resolve_key(f"{tr['prefix']}{tr['id']}")
                    if not tk:
                        continue
                    aid = f"{fk}->{tk}"
                    if aid not in drawn:
                        drawn.add(aid)
                        self._arr(sl, fk, tk, C.ARROW_AWS, label)

            # Source -> Global target (RDS with suffixes)
            for az_suffix, az_frs in fr_by_az.items():
                if az_suffix == "_global":
                    continue
                for fr in az_frs:
                    fk = self._resolve_key(f"{fr['prefix']}{fr['id']}")
                    if not fk:
                        continue
                    for tr in tr_by_az.get("_global", []):
                        if tr['type'] == 'RDS':
                            rds_suffix = "_0" if az_suffix.endswith("a") \
                                else "_1"
                            tk = f"{tr['prefix']}{tr['id']}{rds_suffix}"
                            if tk in self.pos:
                                aid = f"{fk}->{tk}"
                                if aid not in drawn:
                                    drawn.add(aid)
                                    self._arr(sl, fk, tk, C.ARROW_AWS, label)
                        else:
                            tk = self._resolve_key(
                                f"{tr['prefix']}{tr['id']}")
                            if tk:
                                aid = f"{fk}->{tk}"
                                if aid not in drawn:
                                    drawn.add(aid)
                                    self._arr(sl, fk, tk, C.ARROW_AWS, label)

        # ---- WAF -> ALB (filter) ----
        # WAF has explicit relationship to ALB in JSON
        if waf:
            for alb in albs:
                alb_key = f"alb_{alb['id']}"
                if "waf" in self.pos and alb_key in self.pos:
                    self._arr(sl, "waf", alb_key, C.ARROW_AWS, "Filter")

        # ---- Service-level connections ----
        for conn in svc_conns:
            ft = conn.get("from_type", "")
            fi = conn.get("from_id", "")
            tt = conn.get("to_type", "")
            ti = conn.get("to_id", "")
            label = conn.get("label", "")

            fk = self._find_pos_key(ft, fi)
            tk = self._find_pos_key(tt, ti)
            if fk and tk:
                aid = f"{fk}->{tk}"
                if aid not in drawn:
                    drawn.add(aid)
                    self._arr(sl, fk, tk, C.ARROW_GRAY, label)

        # (ASG is shown as annotation, not as icon with arrows)

    def _find_pos_key(self, svc_type, svc_id):
        """Find a position key for a service connection."""
        if svc_type in self.pos:
            return svc_type
        candidates = [f"{svc_type}_{svc_id}", svc_type, svc_id]
        for c in candidates:
            if c in self.pos:
                return c
        return None

    def _group_by_az(self, resources):
        """Group resources by AZ suffix from their name."""
        groups = defaultdict(list)
        for r in resources:
            name = r.get("name", "")
            az_suffix = None
            for part in name.split("-"):
                if len(part) == 2 and part[0].isdigit() and part[1].isalpha():
                    az_suffix = part
                    break
            if az_suffix:
                groups[az_suffix].append(r)
            else:
                groups["_global"].append(r)
        return groups

    def _resolve_key(self, key):
        if key in self.pos:
            return key
        for s in ["_0", "_1"]:
            if key + s in self.pos:
                return key + s
        return None

    # ==========================================================
    # Primitives
    # ==========================================================
    def _box(self, sl, x, y, w, h, fill, border, bw=Pt(1), r=0.015):
        s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                int(x), int(y), int(w), int(h))
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        s.line.color.rgb = border
        s.line.width = bw
        s.adjustments[0] = r
        s.text_frame.clear()
        return s

    def _txt(self, sl, x, y, w, h, text, sz=10, bold=False, color=None,
             align=PP_ALIGN.LEFT):
        tb = sl.shapes.add_textbox(int(x), int(y), int(w), int(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = color or C.TEXT
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        return tb

    def _ilabel(self, sl, x, y, icon, text, sz=8, bold=False, color=None):
        """Icon + inline text label (for box headers)."""
        isz = Inches(0.22)
        if icon and icon in ICONS:
            sl.shapes.add_picture(ICONS[icon], int(x), int(y), isz, isz)
            self._txt(sl, int(x) + isz + Inches(0.04), int(y),
                      Inches(3.5), isz, text, sz, bold, color)
        else:
            self._txt(sl, int(x), int(y), Inches(3.5), Inches(0.22),
                      text, sz, bold, color)

    def _ibox(self, sl, x, y, icon, label, key, nobg=False):
        """Icon + label below. Register bounding box and shape for connectors.

        Stores the picture Shape object in self.shapes[key] so that
        _arr() can use begin_connect/end_connect for true connector binding.
        """
        isz = Inches(0.42)
        tw = Inches(1.2)
        ix = int(x + tw / 2 - isz / 2)

        pic = None
        if icon in ICONS:
            pic = sl.shapes.add_picture(ICONS[icon], ix, int(y), isz, isz)

        self._txt(sl, int(x), int(y) + isz + Inches(0.01),
                  tw, Inches(0.28), label, 6, True, C.TEXT, PP_ALIGN.CENTER)

        # Bounding box based on ICON size only (not text label width)
        icon_cx = int(x + tw / 2)
        icon_cy = int(y + isz / 2)
        icon_hw = int(isz / 2)
        icon_hh = int(isz / 2)
        self.pos[key] = (icon_cx, icon_cy, icon_hw, icon_hh)

        # Store shape reference for connector binding
        if pic is not None:
            self.shapes[key] = pic

    @staticmethod
    def _side_anchor(pos_data, target_cx, target_cy):
        """Calculate arrow anchor point on the side of a bounding box."""
        if len(pos_data) == 4:
            cx, cy, hw, hh = pos_data
        else:
            cx, cy = pos_data
            hw, hh = int(Inches(0.3)), int(Inches(0.3))

        dx = target_cx - cx
        dy = target_cy - cy

        if dx == 0 and dy == 0:
            return cx + hw, cy

        if hw == 0:
            hw = 1
        if hh == 0:
            hh = 1

        if abs(dx) * hh > abs(dy) * hw:
            if dx > 0:
                return cx + hw, cy
            else:
                return cx - hw, cy
        else:
            if dy > 0:
                return cx, cy + hh
            else:
                return cx, cy - hh

    @staticmethod
    def _cxn_idx(pos_data, target_cx, target_cy):
        """Return connection point index (0=top, 1=left, 2=bottom, 3=right).

        Matches the side chosen by _side_anchor for rectangles/pictures.
        """
        if len(pos_data) == 4:
            cx, cy, hw, hh = pos_data
        else:
            cx, cy = pos_data
            hw, hh = int(Inches(0.3)), int(Inches(0.3))

        dx = target_cx - cx
        dy = target_cy - cy

        if dx == 0 and dy == 0:
            return 3  # right

        if hw == 0:
            hw = 1
        if hh == 0:
            hh = 1

        if abs(dx) * hh > abs(dy) * hw:
            return 3 if dx > 0 else 1  # right or left
        else:
            return 2 if dy > 0 else 0  # bottom or top

    def _arr(self, sl, fk, tk, color, label=""):
        """Draw arrow with arrowhead from fk to tk.

        If both shapes are available, uses begin_connect/end_connect for
        true PPTX connector binding (arrows follow icons when moved).
        Falls back to coordinate-based connectors when shapes are missing.
        """
        if fk not in self.pos or tk not in self.pos:
            return

        fp = self.pos[fk]
        tp = self.pos[tk]

        fcx, fcy = fp[0], fp[1]
        tcx, tcy = tp[0], tp[1]

        sx, sy = self._side_anchor(fp, tcx, tcy)
        ex, ey = self._side_anchor(tp, fcx, fcy)

        cn = sl.shapes.add_connector(1, sx, sy, ex, ey)
        cn.line.color.rgb = color
        cn.line.width = Pt(1.5)

        # Bind connector to shapes if both are available
        f_shape = self.shapes.get(fk)
        t_shape = self.shapes.get(tk)
        if f_shape is not None:
            f_idx = self._cxn_idx(fp, tcx, tcy)
            cn.begin_connect(f_shape, f_idx)
        if t_shape is not None:
            t_idx = self._cxn_idx(tp, fcx, fcy)
            cn.end_connect(t_shape, t_idx)

        # Arrowhead
        sp = cn._element.find(qn('p:spPr'))
        ln = sp.find(qn('a:ln'))
        if ln is None:
            ln = etree.SubElement(sp, qn('a:ln'))
        te = ln.find(qn('a:tailEnd'))
        if te is None:
            te = etree.SubElement(ln, qn('a:tailEnd'))
        te.set('type', 'triangle')
        te.set('w', 'med')
        te.set('len', 'med')

        if label:
            dx = ex - sx
            dy = ey - sy
            arr_len = (dx * dx + dy * dy) ** 0.5

            # Skip label on very short arrows to avoid overlap with icons
            if arr_len < Inches(0.5):
                return

            # Place label at 40% along the arrow (closer to source)
            mx = int(sx + dx * 0.4)
            my = int(sy + dy * 0.4)

            # Offset perpendicular to arrow direction
            if abs(dy) > abs(dx):
                mx += int(Inches(0.22))
            else:
                my -= int(Inches(0.18))

            self._txt(sl, mx - int(Inches(0.5)), my - int(Inches(0.1)),
                      Inches(1.1), Inches(0.22), label, 6, True, color,
                      PP_ALIGN.CENTER)

    def _legend(self, sl, x, y):
        """Legend box with arrow color meanings."""
        lw = Inches(2.2)
        lh = Inches(1.3)
        self._box(sl, x, y, lw, lh, C.WHITE, RGBColor(0xCC, 0xCC, 0xCC))
        self._txt(sl, x + Inches(0.1), y + Inches(0.05),
                  Inches(1.5), Inches(0.18), "Legend:", 8, True)

        ly = y + Inches(0.26)
        self._txt(sl, x + Inches(0.1), ly, Inches(0.35), Inches(0.16),
                  "───▶", 7, True, C.ARROW_INET)
        self._txt(sl, x + Inches(0.48), ly, Inches(1.6), Inches(0.16),
                  "Internet traffic", 7)

        ly += Inches(0.2)
        self._txt(sl, x + Inches(0.1), ly, Inches(0.35), Inches(0.16),
                  "───▶", 7, True, C.ARROW_AWS)
        self._txt(sl, x + Inches(0.48), ly, Inches(1.6), Inches(0.16),
                  "AWS internal traffic", 7)

        ly += Inches(0.2)
        self._txt(sl, x + Inches(0.1), ly, Inches(0.35), Inches(0.16),
                  "───▶", 7, True, C.ARROW_PEER)
        self._txt(sl, x + Inches(0.48), ly, Inches(1.6), Inches(0.16),
                  "VPC Peering", 7)

        ly += Inches(0.2)
        self._txt(sl, x + Inches(0.1), ly, Inches(0.35), Inches(0.16),
                  "───▶", 7, True, C.ARROW_GRAY)
        self._txt(sl, x + Inches(0.48), ly, Inches(1.6), Inches(0.16),
                  "Service connection", 7)


# ============================================================
def main():
    if len(sys.argv) < 2:
        print("Usage: python generate_diagram_v2.py <config.json> [--list] [--vpc vpc-id1,vpc-id2,...] [--debug]")
        print()
        print("Options:")
        print("  --list           List all VPCs in the config and exit")
        print("  --vpc id1,id2    Draw specific VPC(s), each on its own slide")
        print("  --debug          Show detailed resource counts per VPC")
        print()
        print("Default: auto-selects the VPC with the most resources")
        sys.exit(1)

    inp = sys.argv[1]
    if not os.path.exists(inp):
        print(f"Error: {inp} not found")
        sys.exit(1)

    out = os.path.join(os.path.dirname(os.path.abspath(inp)),
                       "network_diagram_v2.pptx")
    parser = AWSConfigParser(inp)

    print(f"Parsing: {inp}")
    for rt, items in sorted(parser.by_type.items()):
        print(f"  {rt}: {len(items)}")

    # --list: show all VPCs and exit
    if "--list" in sys.argv:
        dg = DiagramV2(parser)
        vpcs = dg.list_vpcs()
        print(f"\nVPCs found: {len(vpcs)}")
        for v in sorted(vpcs, key=lambda x: -x["score"]):
            default_tag = " (default)" if v.get("is_default") else ""
            print(f"  {v['id']}  {v['name']:30s}  {v['cidr']:18s}  score={v['score']}{default_tag}")
        print(f"\nUsage: python generate_diagram_v2.py {inp} --vpc {vpcs[0]['id']}")
        return

    # --vpc: draw specified VPCs
    vpc_ids = None
    for i, arg in enumerate(sys.argv):
        if arg == "--vpc" and i + 1 < len(sys.argv):
            vpc_ids = [v.strip() for v in sys.argv[i + 1].split(",")]

    # --debug: show resource breakdown per VPC
    debug = "--debug" in sys.argv
    if debug:
        dg = DiagramV2(parser)
        vpcs = dg.list_vpcs()
        for v in sorted(vpcs, key=lambda x: -x["score"]):
            vid = v["id"]
            default_tag = " (default)" if v.get("is_default") else ""
            print(f"\n  VPC: {v['name']} ({vid}) {v['cidr']}{default_tag}")
            subs = parser.get_subnets_for_vpc(vid)
            print(f"    Subnets: {len(subs)}")
            for s in subs:
                print(f"      {s['id']} {s['tier']:10s} {s['az']:20s} {s['cidr']}")
                instances = parser.get_instances_for_subnet(s['id'])
                for inst in instances:
                    print(f"        EC2: {inst['id']} {inst['name']}")
            albs = parser.get_albs_for_vpc(vid)
            print(f"    ALBs: {len(albs)}")
            for a in albs:
                print(f"      {a['id']} {a['name']}")
            rdss = parser.get_rds_for_vpc(vid)
            print(f"    RDS: {len(rdss)}")
            nats = parser.get_nat_gateways_for_vpc(vid)
            print(f"    NAT GW: {len(nats)}")
            igw = parser.get_igw_for_vpc(vid)
            print(f"    IGW: {'yes' if igw else 'no'}")

    print(f"\nGenerating v2 diagram...")
    DiagramV2(parser).generate(out, vpc_ids=vpc_ids)
    print("Done!")


if __name__ == "__main__":
    main()
