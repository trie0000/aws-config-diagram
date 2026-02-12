"""
aws_config_parser.py: AWS Config Snapshot JSON Parser + v1 Diagram Generator

Parses AWS Config Snapshot JSON and extracts VPC, Subnet, EC2, RDS,
ALB, IGW, NAT, Security Group, and other AWS resource information.

AWSConfigParser class is shared by diagram_pptx.py and diagram_excel.py.

Version: 1.1.0
Last Updated: 2026-02-12
"""

import json
import sys
import os
from collections import defaultdict

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ============================================================
# Color Palette
# ============================================================
class Colors:
    # VPC
    VPC_FILL = RGBColor(0xE8, 0xF5, 0xE9)       # light green
    VPC_BORDER = RGBColor(0x2E, 0x7D, 0x32)      # dark green

    # Subnet tiers
    PUBLIC_FILL = RGBColor(0xE3, 0xF2, 0xFD)     # light blue
    PUBLIC_BORDER = RGBColor(0x15, 0x65, 0xC0)    # dark blue
    PRIVATE_FILL = RGBColor(0xFD, 0xF0, 0xE0)    # light orange
    PRIVATE_BORDER = RGBColor(0xE6, 0x51, 0x00)   # dark orange
    ISOLATED_FILL = RGBColor(0xF3, 0xE5, 0xF5)   # light purple
    ISOLATED_BORDER = RGBColor(0x6A, 0x1B, 0x9A)  # dark purple

    # Resources
    IGW_FILL = RGBColor(0xFF, 0xF9, 0xC4)        # yellow
    IGW_BORDER = RGBColor(0xF5, 0x7F, 0x17)
    NAT_FILL = RGBColor(0xFF, 0xEC, 0xB3)
    NAT_BORDER = RGBColor(0xFF, 0x8F, 0x00)
    ALB_FILL = RGBColor(0xBB, 0xDE, 0xFB)
    ALB_BORDER = RGBColor(0x0D, 0x47, 0xA1)
    EC2_FILL = RGBColor(0xFF, 0xCC, 0xBC)
    EC2_BORDER = RGBColor(0xBF, 0x36, 0x0C)
    RDS_FILL = RGBColor(0xC5, 0xCA, 0xE9)
    RDS_BORDER = RGBColor(0x28, 0x35, 0x93)

    # WAF
    WAF_FILL = RGBColor(0xFF, 0xE0, 0xE0)         # light red/pink
    WAF_BORDER = RGBColor(0xC6, 0x28, 0x28)       # dark red
    # S3
    S3_FILL = RGBColor(0xC8, 0xE6, 0xC9)          # light green
    S3_BORDER = RGBColor(0x2E, 0x7D, 0x32)        # dark green

    # Arrows
    ARROW_EXTERNAL = RGBColor(0xFF, 0x00, 0x00)   # red - external access
    ARROW_INTERNAL = RGBColor(0x33, 0x33, 0x33)   # dark gray
    ARROW_PEERING = RGBColor(0x00, 0x96, 0x88)    # teal

    # Alert
    ALERT_RED = RGBColor(0xFF, 0x00, 0x00)
    TEXT_BLACK = RGBColor(0x00, 0x00, 0x00)
    TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    TEXT_GRAY = RGBColor(0x66, 0x66, 0x66)

    # Internet
    INTERNET_FILL = RGBColor(0xE0, 0xE0, 0xE0)
    INTERNET_BORDER = RGBColor(0x61, 0x61, 0x61)

    # Peering
    PEERING_FILL = RGBColor(0xB2, 0xDF, 0xDB)
    PEERING_BORDER = RGBColor(0x00, 0x69, 0x5C)


# ============================================================
# Parse AWS Config Snapshot
# ============================================================
class AWSConfigParser:
    """Parse AWS Config snapshot JSON and extract audit-relevant info."""

    # Resource types relevant for external audit
    AUDIT_RESOURCE_TYPES = {
        # Networking
        "AWS::EC2::VPC",
        "AWS::EC2::Subnet",
        "AWS::EC2::InternetGateway",
        "AWS::EC2::NatGateway",
        "AWS::EC2::Instance",
        "AWS::EC2::SecurityGroup",
        "AWS::EC2::RouteTable",
        "AWS::EC2::VPCPeeringConnection",
        "AWS::ElasticLoadBalancingV2::LoadBalancer",
        "AWS::CloudFront::Distribution",
        "AWS::ApiGateway::RestApi",
        "AWS::ApiGatewayV2::Api",
        "AWS::Route53::HostedZone",
        # Compute
        "AWS::Lambda::Function",
        "AWS::ECS::Cluster",
        "AWS::ECS::Service",
        "AWS::ECS::TaskDefinition",
        "AWS::EKS::Cluster",
        "AWS::AutoScaling::AutoScalingGroup",
        # Database / Storage
        "AWS::RDS::DBInstance",
        "AWS::DynamoDB::Table",
        "AWS::ElastiCache::CacheCluster",
        "AWS::Redshift::Cluster",
        "AWS::S3::Bucket",
        # Messaging
        "AWS::SQS::Queue",
        "AWS::SNS::Topic",
        # Security / Monitoring
        "AWS::WAFv2::WebACL",
        "AWS::KMS::Key",
        "AWS::CloudTrail::Trail",
        "AWS::CloudWatch::Alarm",
        # Other
        "AWS::ElasticBeanstalk::Environment",
    }

    @staticmethod
    def _to_camel(key):
        """Convert PascalCase or mixed key to camelCase.

        CidrBlock -> cidrBlock, VpcId -> vpcId, DBInstanceClass -> dBInstanceClass
        """
        if not key or not key[0].isupper():
            return key
        # Handle leading uppercase run (e.g. DBInstance -> dBInstance)
        i = 0
        while i < len(key) - 1 and key[i].isupper() and key[i + 1].isupper():
            i += 1
        if i == 0:
            return key[0].lower() + key[1:]
        # e.g. "DBInstanceClass" -> i=1, want "dBInstanceClass"
        return key[:i].lower() + key[i:]

    @classmethod
    def _normalize_keys(cls, obj):
        """Recursively convert all dict keys from PascalCase to camelCase."""
        if isinstance(obj, dict):
            return {cls._to_camel(k): cls._normalize_keys(v) for k, v in obj.items()}
        elif isinstance(obj, list):
            return [cls._normalize_keys(item) for item in obj]
        return obj

    def __init__(self, snapshot_path):
        with open(snapshot_path, "r", encoding="utf-8") as f:
            self.data = json.load(f)

        self.items = self.data.get("configurationItems", [])
        self.by_type = defaultdict(list)
        self.by_id = {}

        for item in self.items:
            # Normalize None -> empty dict
            if item.get("configuration") is None:
                item["configuration"] = {}
            if item.get("relationships") is None:
                item["relationships"] = []

            # Parse configuration if it's a JSON string
            cfg = item.get("configuration")
            if isinstance(cfg, str):
                try:
                    item["configuration"] = json.loads(cfg)
                except (json.JSONDecodeError, TypeError):
                    pass

            # Normalize configuration keys: PascalCase -> camelCase
            cfg = item.get("configuration")
            if isinstance(cfg, dict):
                item["configuration"] = self._normalize_keys(cfg)

            # Normalize tags: list-of-dicts -> simple dict
            tags = item.get("tags")
            if isinstance(tags, list):
                tag_dict = {}
                for t in tags:
                    if isinstance(t, dict):
                        k = t.get("key", t.get("Key", ""))
                        v = t.get("value", t.get("Value", ""))
                        if k:
                            tag_dict[k] = v
                item["tags"] = tag_dict
            elif not isinstance(tags, dict):
                item["tags"] = {}

            rt = item.get("resourceType", "")
            rid = item.get("resourceId", "")
            if rt in self.AUDIT_RESOURCE_TYPES:
                self.by_type[rt].append(item)
                self.by_id[rid] = item

    @staticmethod
    def _normalize_ip_ranges(ip_ranges):
        """Handle both formats: ['0.0.0.0/0'] and [{'cidrIp': '0.0.0.0/0'}]"""
        result = []
        for item in ip_ranges:
            if isinstance(item, str):
                result.append({"cidrIp": item})
            else:
                result.append(item)
        return result

    def get_vpcs(self):
        vpcs = []
        for item in self.by_type["AWS::EC2::VPC"]:
            cfg = item.get("configuration", {})
            if not isinstance(cfg, dict):
                cfg = {}
            cidr = cfg.get("cidrBlock", "")
            if not cidr:
                # Try supplementaryConfiguration
                supp = item.get("supplementaryConfiguration", {})
                if isinstance(supp, dict):
                    cidr = supp.get("cidrBlock", "")
                    if not cidr:
                        assocs = supp.get("cidrBlockAssociationSet", [])
                        if isinstance(assocs, str):
                            try:
                                assocs = json.loads(assocs)
                            except Exception:
                                assocs = []
                        for a in assocs:
                            if isinstance(a, dict):
                                cb = a.get("cidrBlock", "")
                                if cb:
                                    cidr = cb
                                    break
            if not cidr:
                rn = item.get("resourceName", "")
                if "/" in rn:
                    cidr = rn
            vpcs.append({
                "id": item["resourceId"],
                "name": item.get("tags", {}).get("Name", item["resourceId"]),
                "cidr": cidr,
                "region": item.get("awsRegion", ""),
                "is_default": cfg.get("isDefault", False),
            })
        return vpcs

    def _build_subnet_tier_map(self, vpc_id):
        """Build subnet->tier map using route table analysis + heuristics.

        Primary: Route Table routes analysis
        - Public: route table has 0.0.0.0/0 -> IGW
        - Private: route table has 0.0.0.0/0 -> NAT GW (or other)
        - Isolated: route table has NO 0.0.0.0/0 route

        Fallback (when Route Table configuration is empty):
        - NAT Gateway subnet -> Public (NAT is placed in Public subnet)
        - ALB subnet -> Public (internet-facing ALB is in Public subnet)
        - RDS subnet -> Isolated (DB subnets are isolated)
        - mapPublicIpOnLaunch -> Public
        """
        subnet_tier = {}

        # --- Source 1: Route Table analysis ---
        for item in self.by_type["AWS::EC2::RouteTable"]:
            cfg = item.get("configuration", {})
            if not isinstance(cfg, dict):
                cfg = {}
            rt_vpc = cfg.get("vpcId", "")
            if not rt_vpc:
                rt_vpc = self._get_related_vpc(item)
            if rt_vpc != vpc_id:
                continue
            routes = cfg.get("routes", cfg.get("routeSet", []))
            has_igw = False
            has_nat = False
            has_default = False
            for r in routes:
                dest = r.get("destinationCidrBlock",
                             r.get("destinationCidr", ""))
                if dest == "0.0.0.0/0":
                    has_default = True
                    gw = r.get("gatewayId", "")
                    nat = r.get("natGatewayId", "")
                    if gw and gw.startswith("igw-"):
                        has_igw = True
                    elif nat and nat.startswith("nat-"):
                        has_nat = True

            if has_igw:
                tier = "Public"
            elif has_nat or has_default:
                tier = "Private"
            else:
                tier = "Isolated"

            # Map associated subnets
            assocs = cfg.get("associations",
                             cfg.get("routeTableAssociationSet", []))
            for a in assocs:
                sid = a.get("subnetId", "")
                if sid:
                    subnet_tier[sid] = tier
                elif a.get("main", False):
                    # Main route table — serves as default for unassociated subs
                    subnet_tier["_main"] = tier

        # If Route Table analysis produced results, return early
        real_entries = {k: v for k, v in subnet_tier.items() if k != "_main"}
        if real_entries:
            return subnet_tier

        # --- Source 2: Heuristic fallback (Route Table data unavailable) ---
        s2v = self._build_subnet_vpc_map()
        vpc_subnet_ids = {sid for sid, vid in s2v.items() if vid == vpc_id}

        # Try direct subnet mapping from service configurations
        # NAT Gateway subnet -> Public (NAT GW resides in a Public subnet)
        for item in self.by_type["AWS::EC2::NatGateway"]:
            cfg = item.get("configuration", {})
            if not isinstance(cfg, dict):
                cfg = {}
            sub = cfg.get("subnetId", "")
            if sub and sub in vpc_subnet_ids and sub not in subnet_tier:
                subnet_tier[sub] = "Public"

        # ALB (internet-facing) subnets -> Public
        for item in self.by_type["AWS::ElasticLoadBalancingV2::LoadBalancer"]:
            cfg = item.get("configuration", {})
            if not isinstance(cfg, dict):
                cfg = {}
            scheme = cfg.get("scheme", "")
            if scheme == "internet-facing":
                for az in cfg.get("availabilityZones", []):
                    sub = az.get("subnetId", "")
                    if sub and sub in vpc_subnet_ids and sub not in subnet_tier:
                        subnet_tier[sub] = "Public"

        # RDS subnets -> Isolated
        for item in self.by_type["AWS::RDS::DBInstance"]:
            cfg = item.get("configuration", {})
            if not isinstance(cfg, dict):
                cfg = {}
            sg_group = cfg.get("dBSubnetGroup", cfg.get("dbSubnetGroup", {}))
            if not isinstance(sg_group, dict):
                sg_group = {}
            for s in sg_group.get("subnets", []):
                sub = s.get("subnetIdentifier", "")
                if sub and sub in vpc_subnet_ids and sub not in subnet_tier:
                    subnet_tier[sub] = "Isolated"

        # If heuristics assigned any tiers, set default for remaining and return
        real_entries2 = {k: v for k, v in subnet_tier.items() if k != "_main"}
        if real_entries2:
            if "_main" not in subnet_tier:
                subnet_tier["_main"] = "Private"
            return subnet_tier

        # --- Source 3: Last resort — auto-distribute tiers by AZ ---
        # When ALL configuration data is empty (Route Tables, NAT subnetId,
        # RDS subnet_ids, ALB subnets, etc.), distribute subnets into
        # Public/Private/Isolated based on service existence and AZ grouping.
        has_igw = bool(self.by_type["AWS::EC2::InternetGateway"])
        has_nat = bool(self.by_type["AWS::EC2::NatGateway"])
        has_rds = bool(self.by_type["AWS::RDS::DBInstance"])

        # Group VPC subnets by AZ
        az_subnets = defaultdict(list)
        for item in self.by_type["AWS::EC2::Subnet"]:
            sid = item["resourceId"]
            if sid not in vpc_subnet_ids:
                continue
            az = ""
            cfg_s = item.get("configuration", {})
            if isinstance(cfg_s, dict):
                az = cfg_s.get("availabilityZone", "")
            if not az:
                az = item.get("availabilityZone", "")
            az_subnets[az or "unknown"].append(sid)

        # Find subnets that have EC2 instances (these are Private/app subnets)
        ec2_subnets = set()
        for item in self.by_type["AWS::EC2::Instance"]:
            cfg_e = item.get("configuration", {})
            if isinstance(cfg_e, dict):
                sub = cfg_e.get("subnetId", "")
                if sub and sub in vpc_subnet_ids:
                    ec2_subnets.add(sub)

        # Distribute: for each AZ, assign tiers to subnets
        for az, sids in az_subnets.items():
            if len(sids) == 1:
                # Only 1 subnet in this AZ: make it Private (most useful)
                subnet_tier[sids[0]] = "Private"
            elif len(sids) == 2:
                # 2 subnets: Public + Private
                if has_igw or has_nat:
                    # The one WITHOUT EC2 is likely Public
                    non_ec2 = [s for s in sids if s not in ec2_subnets]
                    with_ec2 = [s for s in sids if s in ec2_subnets]
                    if non_ec2 and with_ec2:
                        subnet_tier[non_ec2[0]] = "Public"
                        subnet_tier[with_ec2[0]] = "Private"
                    else:
                        subnet_tier[sids[0]] = "Public"
                        subnet_tier[sids[1]] = "Private"
                else:
                    subnet_tier[sids[0]] = "Public"
                    subnet_tier[sids[1]] = "Private"
            else:
                # 3+ subnets: Public + Private + Isolated
                non_ec2 = [s for s in sids if s not in ec2_subnets]
                with_ec2 = [s for s in sids if s in ec2_subnets]
                assigned = set()

                # First non-EC2 subnet -> Public
                if non_ec2:
                    subnet_tier[non_ec2[0]] = "Public"
                    assigned.add(non_ec2[0])
                elif sids:
                    subnet_tier[sids[0]] = "Public"
                    assigned.add(sids[0])

                # If RDS exists, assign one subnet as Isolated
                if has_rds:
                    for s in non_ec2:
                        if s not in assigned:
                            subnet_tier[s] = "Isolated"
                            assigned.add(s)
                            break
                    else:
                        # No remaining non-EC2 subnet; use last one
                        for s in reversed(sids):
                            if s not in assigned:
                                subnet_tier[s] = "Isolated"
                                assigned.add(s)
                                break

                # Remaining -> Private
                for s in sids:
                    if s not in assigned:
                        subnet_tier[s] = "Private"

        if "_main" not in subnet_tier:
            subnet_tier["_main"] = "Private"

        return subnet_tier

    @staticmethod
    def _get_related_vpc(item):
        """Extract VPC ID from relationships when configuration is empty."""
        for rel in item.get("relationships", []):
            if rel.get("resourceType") == "AWS::EC2::VPC":
                return rel.get("resourceId", "")
        return ""

    def _build_subnet_vpc_map(self):
        """Build subnet_id -> vpc_id map from all available sources.

        When Subnet configuration is empty, we reverse-engineer the mapping
        from EC2 instances, ALBs, NAT Gateways, etc. that DO have both
        subnetId and vpcId in their configuration.
        """
        s2v = {}

        # Source 1: Subnet configuration itself
        for item in self.by_type["AWS::EC2::Subnet"]:
            cfg = item.get("configuration", {})
            if isinstance(cfg, dict) and cfg.get("vpcId"):
                s2v[item["resourceId"]] = cfg["vpcId"]

        # Source 2: Subnet relationships
        for item in self.by_type["AWS::EC2::Subnet"]:
            sid = item["resourceId"]
            if sid not in s2v:
                vpc = self._get_related_vpc(item)
                if vpc:
                    s2v[sid] = vpc

        # Source 3: EC2 instances (subnetId + vpcId in config)
        for item in self.by_type["AWS::EC2::Instance"]:
            cfg = item.get("configuration", {})
            if isinstance(cfg, dict):
                sub = cfg.get("subnetId", "")
                vpc = cfg.get("vpcId", "")
                if sub and vpc and sub not in s2v:
                    s2v[sub] = vpc

        # Source 4: NAT Gateways
        for item in self.by_type["AWS::EC2::NatGateway"]:
            cfg = item.get("configuration", {})
            if isinstance(cfg, dict):
                sub = cfg.get("subnetId", "")
                vpc = cfg.get("vpcId", "")
                if sub and vpc and sub not in s2v:
                    s2v[sub] = vpc

        # Source 5: ALB availability zones
        for item in self.by_type["AWS::ElasticLoadBalancingV2::LoadBalancer"]:
            cfg = item.get("configuration", {})
            if isinstance(cfg, dict):
                vpc = cfg.get("vpcId", "")
                if vpc:
                    for az in cfg.get("availabilityZones", []):
                        sub = az.get("subnetId", "")
                        if sub and sub not in s2v:
                            s2v[sub] = vpc

        # Source 6: Route Table associations (subnetId in associations)
        for item in self.by_type["AWS::EC2::RouteTable"]:
            cfg = item.get("configuration", {})
            if not isinstance(cfg, dict):
                continue
            rt_vpc = cfg.get("vpcId", "")
            if not rt_vpc:
                rt_vpc = self._get_related_vpc(item)
            if not rt_vpc:
                continue
            assocs = cfg.get("associations",
                             cfg.get("routeTableAssociationSet", []))
            for a in assocs:
                sid = a.get("subnetId", "")
                if sid and sid not in s2v:
                    s2v[sid] = rt_vpc

        return s2v

    def _get_primary_vpc_id(self):
        """Get the primary (non-default, most-resourced) VPC ID.

        Used as last-resort fallback when resource-to-VPC mapping fails
        for all other methods.
        """
        s2v = self._build_subnet_vpc_map()
        # Count resources per VPC from subnet map
        vpc_counts = defaultdict(int)
        for vid in s2v.values():
            vpc_counts[vid] += 1
        # Also count EC2 instances per VPC
        for item in self.by_type["AWS::EC2::Instance"]:
            cfg = item.get("configuration", {})
            if isinstance(cfg, dict) and cfg.get("vpcId"):
                vpc_counts[cfg["vpcId"]] += 10

        if vpc_counts:
            return max(vpc_counts, key=vpc_counts.get)

        # Fallback: first non-default VPC, or just first VPC
        for item in self.by_type["AWS::EC2::VPC"]:
            cfg = item.get("configuration", {})
            if not cfg.get("isDefault", False):
                return item["resourceId"]
        if self.by_type["AWS::EC2::VPC"]:
            return self.by_type["AWS::EC2::VPC"][0]["resourceId"]
        return ""

    def _build_igw_vpc_map(self):
        """Build igw_id -> vpc_id map from Route Tables.

        When IGW configuration/relationships are empty, we reverse-engineer
        the mapping from Route Tables that reference the IGW in their routes.
        Last resort: assign unmatched IGWs to the primary VPC.
        """
        igw2vpc = {}

        # Source 1: IGW configuration.attachments
        for item in self.by_type["AWS::EC2::InternetGateway"]:
            cfg = item.get("configuration", {})
            if not isinstance(cfg, dict):
                cfg = {}
            for att in cfg.get("attachments", []):
                vid = att.get("vpcId", "")
                if vid:
                    igw2vpc[item["resourceId"]] = vid
            # Source 2: relationships
            if item["resourceId"] not in igw2vpc:
                rel_vpc = self._get_related_vpc(item)
                if rel_vpc:
                    igw2vpc[item["resourceId"]] = rel_vpc

        # Source 3: Route Tables (routes contain gatewayId = igw-xxx)
        for item in self.by_type["AWS::EC2::RouteTable"]:
            cfg = item.get("configuration", {})
            if not isinstance(cfg, dict):
                continue
            rt_vpc = cfg.get("vpcId", "")
            if not rt_vpc:
                rt_vpc = self._get_related_vpc(item)
            if not rt_vpc:
                continue
            routes = cfg.get("routes", cfg.get("routeSet", []))
            for r in routes:
                gw = r.get("gatewayId", "")
                if gw and gw.startswith("igw-") and gw not in igw2vpc:
                    igw2vpc[gw] = rt_vpc

        # Source 4 (last resort): Assign unmatched IGWs to primary VPC
        # IGW is typically 1:1 with VPC; if all other methods fail, assume
        # unmatched IGWs belong to the VPC with the most resources
        unmatched = [item["resourceId"] for item in self.by_type["AWS::EC2::InternetGateway"]
                     if item["resourceId"] not in igw2vpc]
        if unmatched:
            primary_vpc = self._get_primary_vpc_id()
            if primary_vpc:
                # Only assign if the primary VPC doesn't already have an IGW
                vpcs_with_igw = set(igw2vpc.values())
                if primary_vpc not in vpcs_with_igw:
                    # Assign the first unmatched IGW to primary VPC
                    igw2vpc[unmatched[0]] = primary_vpc

        return igw2vpc

    def _build_nat_vpc_map(self):
        """Build nat_id -> vpc_id map from Route Tables and NAT config.

        When NAT configuration/relationships are empty, we reverse-engineer
        the mapping from Route Tables that reference the NAT in their routes.
        Last resort: assign unmatched NATs to the primary VPC.
        """
        nat2vpc = {}

        # Source 1: NAT configuration.vpcId
        for item in self.by_type["AWS::EC2::NatGateway"]:
            cfg = item.get("configuration", {})
            if not isinstance(cfg, dict):
                cfg = {}
            vid = cfg.get("vpcId", "")
            if vid:
                nat2vpc[item["resourceId"]] = vid
            elif self._get_related_vpc(item):
                nat2vpc[item["resourceId"]] = self._get_related_vpc(item)

        # Source 2: Route Tables (routes contain natGatewayId = nat-xxx)
        for item in self.by_type["AWS::EC2::RouteTable"]:
            cfg = item.get("configuration", {})
            if not isinstance(cfg, dict):
                continue
            rt_vpc = cfg.get("vpcId", "")
            if not rt_vpc:
                rt_vpc = self._get_related_vpc(item)
            if not rt_vpc:
                continue
            routes = cfg.get("routes", cfg.get("routeSet", []))
            for r in routes:
                nat = r.get("natGatewayId", "")
                if nat and nat.startswith("nat-") and nat not in nat2vpc:
                    nat2vpc[nat] = rt_vpc

        # Source 3 (last resort): Assign unmatched NATs to primary VPC
        unmatched = [item["resourceId"] for item in self.by_type["AWS::EC2::NatGateway"]
                     if item["resourceId"] not in nat2vpc]
        if unmatched:
            primary_vpc = self._get_primary_vpc_id()
            if primary_vpc:
                for nat_id in unmatched:
                    nat2vpc[nat_id] = primary_vpc

        return nat2vpc

    def _build_rds_vpc_map(self):
        """Build rds_id -> vpc_id map by matching RDS subnet IDs to known subnets.

        When RDS dBSubnetGroup.vpcId and relationships are empty, we
        reverse-engineer by finding which VPC the RDS's subnets belong to.
        """
        rds2vpc = {}
        s2v = self._build_subnet_vpc_map()

        for item in self.by_type["AWS::RDS::DBInstance"]:
            cfg = item.get("configuration", {})
            if not isinstance(cfg, dict):
                cfg = {}
            rid = item["resourceId"]

            # Source 1: dBSubnetGroup.vpcId
            sg_group = cfg.get("dBSubnetGroup", cfg.get("dbSubnetGroup", {}))
            if not isinstance(sg_group, dict):
                sg_group = {}
            vid = sg_group.get("vpcId", "")
            if vid:
                rds2vpc[rid] = vid
                continue

            # Source 2: relationships
            rel_vpc = self._get_related_vpc(item)
            if rel_vpc:
                rds2vpc[rid] = rel_vpc
                continue

            # Source 3: Match subnets in dBSubnetGroup to known subnet→VPC map
            for s in sg_group.get("subnets", []):
                sub_id = s.get("subnetIdentifier", "")
                if sub_id and sub_id in s2v:
                    rds2vpc[rid] = s2v[sub_id]
                    break

        # Source 4 (last resort): assign unmapped RDS to primary VPC
        if rds2vpc:
            primary = self._get_primary_vpc_id()
        else:
            primary = ""
        for item in self.by_type["AWS::RDS::DBInstance"]:
            rid = item["resourceId"]
            if rid not in rds2vpc and primary:
                rds2vpc[rid] = primary

        return rds2vpc

    def _build_subnet_cidr_map(self):
        """Build subnet -> CIDR map from NetworkInterface configurations.

        When Subnet configuration is ResourceNotRecorded, we can infer
        the subnet CIDR from NetworkInterface privateIpAddress + subnetId.
        This is a heuristic — we collect all private IPs per subnet and
        try to guess the CIDR from the common prefix.
        """
        subnet_ips = defaultdict(list)
        for item in self.by_type.get("AWS::EC2::NetworkInterface", []):
            cfg = item.get("configuration", {})
            if not isinstance(cfg, dict):
                cfg = {}
            sub_id = cfg.get("subnetId", "")
            if not sub_id:
                continue
            # Collect private IPs
            priv_addrs = cfg.get("privateIpAddresses", [])
            for pa in priv_addrs:
                if isinstance(pa, dict):
                    ip = pa.get("privateIpAddress", "")
                    if ip:
                        subnet_ips[sub_id].append(ip)
            # Also try top-level privateIpAddress
            pip = cfg.get("privateIpAddress", "")
            if pip and pip not in subnet_ips[sub_id]:
                subnet_ips[sub_id].append(pip)

        # Also collect from EC2 instances
        for item in self.by_type.get("AWS::EC2::Instance", []):
            cfg = item.get("configuration", {})
            if not isinstance(cfg, dict):
                cfg = {}
            sub_id = cfg.get("subnetId", "")
            pip = cfg.get("privateIpAddress", "")
            if sub_id and pip:
                if pip not in subnet_ips[sub_id]:
                    subnet_ips[sub_id].append(pip)

        # Try to guess CIDR from collected IPs
        cidr_map = {}
        for sub_id, ips in subnet_ips.items():
            if ips:
                # Use first IP as representative; guess /24 as common default
                ip = ips[0]
                parts = ip.split(".")
                if len(parts) == 4:
                    cidr_map[sub_id] = f"{parts[0]}.{parts[1]}.{parts[2]}.0/24"
        return cidr_map

    def get_subnets_for_vpc(self, vpc_id):
        tier_map = self._build_subnet_tier_map(vpc_id)
        default_tier = tier_map.get("_main", "Private")
        subnet_vpc_map = self._build_subnet_vpc_map()
        subnet_cidr_map = self._build_subnet_cidr_map()
        subnets = []
        for item in self.by_type["AWS::EC2::Subnet"]:
            cfg = item.get("configuration", {})
            if not isinstance(cfg, dict):
                cfg = {}
            sid = item["resourceId"]

            # VPC ID: configuration > relationships > reverse-engineered map
            item_vpc = cfg.get("vpcId", "")
            if not item_vpc:
                item_vpc = self._get_related_vpc(item)
            if not item_vpc:
                item_vpc = subnet_vpc_map.get(sid, "")

            if item_vpc != vpc_id:
                continue

            tags = item.get("tags", {})
            sid = item["resourceId"]

            # AZ: try configuration, fallback to top-level field
            az = cfg.get("availabilityZone", "")
            if not az:
                az = item.get("availabilityZone", "")

            # CIDR: try multiple sources
            cidr = cfg.get("cidrBlock", "")
            if not cidr:
                # Try supplementaryConfiguration
                supp = item.get("supplementaryConfiguration", {})
                if isinstance(supp, dict):
                    # Some snapshots store CIDR in supplementaryConfiguration
                    cidr = supp.get("cidrBlock", "")
                    if not cidr:
                        # cidrBlockAssociationSet
                        assocs = supp.get("cidrBlockAssociationSet", [])
                        if isinstance(assocs, str):
                            try:
                                assocs = json.loads(assocs)
                            except Exception:
                                assocs = []
                        for a in assocs:
                            if isinstance(a, dict):
                                cb = a.get("cidrBlock", "")
                                if cb:
                                    cidr = cb
                                    break
            if not cidr:
                # resourceName sometimes contains CIDR
                rn = item.get("resourceName", "")
                if "/" in rn:
                    cidr = rn
            if not cidr:
                # Fallback: infer from NetworkInterface/EC2 IPs
                cidr = subnet_cidr_map.get(sid, "")

            # Priority: explicit Tier tag > route table > name hint > default
            tier = tags.get("Tier", "")
            if not tier:
                tier = tier_map.get(sid, "")
            if not tier:
                # Heuristic: infer from Name tag
                name = tags.get("Name", "").lower()
                if "public" in name:
                    tier = "Public"
                elif "isolated" in name or "db" in name or "data" in name:
                    tier = "Isolated"
                elif "private" in name:
                    tier = "Private"
            if not tier:
                # mapPublicIpOnLaunch hint
                if cfg.get("mapPublicIpOnLaunch"):
                    tier = "Public"
            if not tier:
                tier = default_tier

            subnets.append({
                "id": sid,
                "name": tags.get("Name", sid),
                "cidr": cidr,
                "az": az,
                "tier": tier,
            })
        return subnets

    def get_igw_for_vpc(self, vpc_id):
        igw_vpc_map = self._build_igw_vpc_map()
        for item in self.by_type["AWS::EC2::InternetGateway"]:
            rid = item["resourceId"]
            # Use comprehensive reverse map (config + relationships + route tables)
            if igw_vpc_map.get(rid) == vpc_id:
                return {
                    "id": rid,
                    "name": item.get("tags", {}).get("Name", rid),
                }
        return None

    def get_nat_gateways_for_vpc(self, vpc_id):
        nat_vpc_map = self._build_nat_vpc_map()
        nats = []
        for item in self.by_type["AWS::EC2::NatGateway"]:
            cfg = item.get("configuration", {})
            if not isinstance(cfg, dict):
                cfg = {}
            rid = item["resourceId"]
            # Use comprehensive reverse map (config + relationships + route tables)
            if nat_vpc_map.get(rid) == vpc_id:
                addrs = cfg.get("natGatewayAddresses", [])
                public_ip = addrs[0].get("publicIp", "") if addrs else ""
                nats.append({
                    "id": rid,
                    "name": item.get("tags", {}).get("Name", rid),
                    "subnet_id": cfg.get("subnetId", ""),
                    "public_ip": public_ip,
                })
        return nats

    def get_instances_for_subnet(self, subnet_id):
        instances = []
        for item in self.by_type["AWS::EC2::Instance"]:
            cfg = item.get("configuration", {})
            if not isinstance(cfg, dict):
                cfg = {}
            # Primary: configuration.subnetId
            matched = cfg.get("subnetId") == subnet_id
            # Fallback: relationships → Subnet reference
            if not matched:
                for rel in item.get("relationships", []):
                    if (rel.get("resourceType") == "AWS::EC2::Subnet"
                            and rel.get("resourceId") == subnet_id):
                        matched = True
                        break
            if matched:
                tags = item.get("tags", {})
                instances.append({
                    "id": item["resourceId"],
                    "name": tags.get("Name", item["resourceId"]),
                    "type": cfg.get("instanceType", ""),
                    "private_ip": cfg.get("privateIpAddress", ""),
                    "public_ip": cfg.get("publicIpAddress"),
                    "role": tags.get("Role", ""),
                    "sg_ids": [sg.get("groupId", "") for sg in cfg.get("securityGroups", [])],
                })
        return instances

    def get_albs_for_vpc(self, vpc_id):
        s2v = self._build_subnet_vpc_map()
        albs = []
        for item in self.by_type["AWS::ElasticLoadBalancingV2::LoadBalancer"]:
            cfg = item.get("configuration", {})
            if not isinstance(cfg, dict):
                cfg = {}
            item_vpc = cfg.get("vpcId", "")
            if not item_vpc:
                item_vpc = self._get_related_vpc(item)
            # Fallback: match ALB's subnets to known subnet→VPC map
            if not item_vpc:
                for az in cfg.get("availabilityZones", []):
                    sub = az.get("subnetId", "")
                    if sub and sub in s2v:
                        item_vpc = s2v[sub]
                        break
            if item_vpc == vpc_id:
                subnet_ids = []
                for az in cfg.get("availabilityZones", []):
                    sid = az.get("subnetId", "")
                    if sid:
                        subnet_ids.append(sid)
                albs.append({
                    "id": item["resourceId"],
                    "name": cfg.get("loadBalancerName", item.get("tags", {}).get("Name", "")),
                    "scheme": cfg.get("scheme", ""),
                    "type": cfg.get("type", ""),
                    "dns": cfg.get("dNSName", ""),
                    "subnet_ids": subnet_ids,
                    "sg_ids": cfg.get("securityGroups", []),
                })
        return albs

    def get_rds_for_vpc(self, vpc_id):
        rds_vpc_map = self._build_rds_vpc_map()
        dbs = []
        for item in self.by_type["AWS::RDS::DBInstance"]:
            cfg = item.get("configuration", {})
            if not isinstance(cfg, dict):
                cfg = {}
            rid = item["resourceId"]
            # Use comprehensive reverse map (subnetGroup + relationships + subnet matching)
            if rds_vpc_map.get(rid) == vpc_id:
                sg_group = cfg.get("dBSubnetGroup", cfg.get("dbSubnetGroup", {}))
                if not isinstance(sg_group, dict):
                    sg_group = {}
                subnet_ids = [s.get("subnetIdentifier", "") for s in sg_group.get("subnets", [])]
                dbs.append({
                    "id": rid,
                    "name": cfg.get("dBInstanceIdentifier", ""),
                    "engine": cfg.get("engine", ""),
                    "instance_class": cfg.get("dBInstanceClass", ""),
                    "port": cfg.get("endpoint", {}).get("port", ""),
                    "multi_az": cfg.get("multiAZ", False),
                    "publicly_accessible": cfg.get("publiclyAccessible", False),
                    "subnet_ids": subnet_ids,
                    "sg_ids": [sg.get("vpcSecurityGroupId", "") for sg in cfg.get("vpcSecurityGroups", [])],
                })
        return dbs

    def get_security_groups_for_vpc(self, vpc_id):
        sgs = []
        for item in self.by_type["AWS::EC2::SecurityGroup"]:
            cfg = item.get("configuration", {})
            if not isinstance(cfg, dict):
                cfg = {}
            # Primary: configuration.vpcId
            item_vpc = cfg.get("vpcId", "")
            # Fallback: relationships → VPC reference
            if not item_vpc:
                item_vpc = self._get_related_vpc(item)
            if item_vpc == vpc_id:
                sgs.append({
                    "id": cfg.get("groupId", item["resourceId"]),
                    "name": cfg.get("groupName", ""),
                    "description": cfg.get("description", ""),
                    "ingress": cfg.get("ipPermissions", []),
                    "egress": cfg.get("ipPermissionsEgress", []),
                })
        return sgs

    def get_peering_connections(self):
        peerings = []
        for item in self.by_type["AWS::EC2::VPCPeeringConnection"]:
            cfg = item.get("configuration", {})
            peerings.append({
                "id": item["resourceId"],
                "name": item.get("tags", {}).get("Name", item["resourceId"]),
                "accepter_vpc": cfg.get("accepterVpcInfo", {}).get("vpcId", ""),
                "requester_vpc": cfg.get("requesterVpcInfo", {}).get("vpcId", ""),
                "accepter_cidr": cfg.get("accepterVpcInfo", {}).get("cidrBlock", ""),
                "requester_cidr": cfg.get("requesterVpcInfo", {}).get("cidrBlock", ""),
            })
        return peerings

    def get_external_sg_rules(self, vpc_id):
        """Extract SG rules that allow 0.0.0.0/0 (external access) - audit critical."""
        external_rules = []
        for sg in self.get_security_groups_for_vpc(vpc_id):
            for rule in sg.get("ingress", []):
                for ip_range in self._normalize_ip_ranges(rule.get("ipRanges", [])):
                    if ip_range.get("cidrIp") == "0.0.0.0/0":
                        external_rules.append({
                            "sg_id": sg["id"],
                            "sg_name": sg["name"],
                            "direction": "INBOUND",
                            "protocol": rule.get("ipProtocol", ""),
                            "from_port": rule.get("fromPort", ""),
                            "to_port": rule.get("toPort", ""),
                            "source": "0.0.0.0/0",
                            "description": ip_range.get("description", ""),
                        })
        return external_rules

    def get_sg_connections(self):
        """Extract SG-to-SG references to draw internal traffic flows."""
        connections = []
        all_sgs = self.by_type["AWS::EC2::SecurityGroup"]
        for item in all_sgs:
            cfg = item.get("configuration", {})
            sg_id = cfg.get("groupId", item["resourceId"])
            sg_name = cfg.get("groupName", "")
            # Ingress: who can talk TO this SG
            for rule in cfg.get("ipPermissions", []):
                for pair in rule.get("userIdGroupPairs", []):
                    connections.append({
                        "from_sg": pair["groupId"],
                        "to_sg": sg_id,
                        "to_sg_name": sg_name,
                        "port": rule.get("fromPort", ""),
                        "protocol": rule.get("ipProtocol", ""),
                        "description": pair.get("description", ""),
                    })
        return connections

    def get_internet_facing_sgs(self):
        """Return SGs that allow inbound from 0.0.0.0/0 (internet-facing).

        Returns list of dicts:
          { "sg_id": str, "port": int|str, "protocol": str }
        """
        results = []
        for item in self.by_type["AWS::EC2::SecurityGroup"]:
            cfg = item.get("configuration", {})
            sg_id = cfg.get("groupId", item["resourceId"])
            for rule in cfg.get("ipPermissions", []):
                ip_ranges = rule.get("ipRanges", [])
                ip_ranges = self._normalize_ip_ranges(ip_ranges)
                for ip_range in ip_ranges:
                    cidr = ip_range.get("cidrIp", "")
                    if cidr == "0.0.0.0/0":
                        results.append({
                            "sg_id": sg_id,
                            "port": rule.get("fromPort", ""),
                            "protocol": rule.get("ipProtocol", ""),
                        })
        return results

    def get_waf_for_alb(self, alb_id_or_arn):
        """Get WAF WebACL associated with an ALB.

        Accepts ALB id (short) or full ARN.  Checks both
        'associatedResources' and '_associated_resources' keys,
        and also tries partial match for short IDs.
        """
        for item in self.by_type["AWS::WAFv2::WebACL"]:
            cfg = item.get("configuration", {})
            associated = (cfg.get("associatedResources", [])
                          + cfg.get("_associated_resources", []))
            # Check exact match or partial (short ALB id in ARN)
            for res_arn in associated:
                if (alb_id_or_arn == res_arn
                        or alb_id_or_arn in res_arn):
                    rules = cfg.get("rules", [])
                    rule_names = [r.get("name", "") for r in rules[:3]]
                    return {
                        "id": item["resourceId"],
                        "name": cfg.get("name", ""),
                        "rules_summary": rule_names,
                        "rule_count": len(rules),
                    }
        return None

    def get_s3_buckets(self):
        """Get S3 buckets."""
        buckets = []
        for item in self.by_type["AWS::S3::Bucket"]:
            cfg = item.get("configuration", {})
            pub_block = cfg.get("publicAccessBlockConfiguration", {})
            is_public_blocked = all([
                pub_block.get("blockPublicAcls", False),
                pub_block.get("blockPublicPolicy", False),
            ])
            buckets.append({
                "id": item["resourceId"],
                "name": cfg.get("name", item["resourceId"]),
                "encrypted": bool(cfg.get("bucketEncryptionConfiguration")),
                "versioning": cfg.get("versioningConfiguration", {}).get("status", "Disabled"),
                "public_blocked": is_public_blocked,
            })
        return buckets

    # ---- New service parsers ----

    def get_lambda_functions(self):
        funcs = []
        for item in self.by_type["AWS::Lambda::Function"]:
            cfg = item.get("configuration", {})
            vpc_cfg = cfg.get("vpcConfig", {})
            funcs.append({
                "id": item["resourceId"],
                "name": cfg.get("functionName", item.get("tags", {}).get("Name", item["resourceId"])),
                "runtime": cfg.get("runtime", ""),
                "memory": cfg.get("memorySize", ""),
                "timeout": cfg.get("timeout", ""),
                "vpc_subnet_ids": vpc_cfg.get("subnetIds", []),
                "vpc_sg_ids": vpc_cfg.get("securityGroupIds", []),
                "in_vpc": bool(vpc_cfg.get("subnetIds")),
            })
        return funcs

    def get_ecs_clusters(self):
        clusters = []
        for item in self.by_type["AWS::ECS::Cluster"]:
            cfg = item.get("configuration", {})
            clusters.append({
                "id": item["resourceId"],
                "name": cfg.get("clusterName", item.get("tags", {}).get("Name", item["resourceId"])),
            })
        return clusters

    def get_ecs_services(self):
        services = []
        for item in self.by_type["AWS::ECS::Service"]:
            cfg = item.get("configuration", {})
            net_cfg = cfg.get("networkConfiguration", {}).get("awsvpcConfiguration", {})
            services.append({
                "id": item["resourceId"],
                "name": cfg.get("serviceName", item.get("tags", {}).get("Name", item["resourceId"])),
                "cluster_arn": cfg.get("clusterArn", ""),
                "launch_type": cfg.get("launchType", ""),
                "desired_count": cfg.get("desiredCount", 0),
                "subnet_ids": net_cfg.get("subnets", []),
                "sg_ids": net_cfg.get("securityGroups", []),
            })
        return services

    def get_eks_clusters(self):
        clusters = []
        for item in self.by_type["AWS::EKS::Cluster"]:
            cfg = item.get("configuration", {})
            vpc_cfg = cfg.get("resourcesVpcConfig", {})
            clusters.append({
                "id": item["resourceId"],
                "name": cfg.get("name", item.get("tags", {}).get("Name", item["resourceId"])),
                "version": cfg.get("version", ""),
                "subnet_ids": vpc_cfg.get("subnetIds", []),
                "sg_ids": vpc_cfg.get("securityGroupIds", []),
            })
        return clusters

    def get_autoscaling_groups(self):
        asgs = []
        for item in self.by_type["AWS::AutoScaling::AutoScalingGroup"]:
            cfg = item.get("configuration", {})
            subnet_str = cfg.get("vPCZoneIdentifier", "")
            asgs.append({
                "id": item["resourceId"],
                "name": cfg.get("autoScalingGroupName", item.get("tags", {}).get("Name", item["resourceId"])),
                "min_size": cfg.get("minSize", 0),
                "max_size": cfg.get("maxSize", 0),
                "desired": cfg.get("desiredCapacity", 0),
                "subnet_ids": [s for s in subnet_str.split(",") if s] if subnet_str else [],
            })
        return asgs

    def get_cloudfront_distributions(self):
        dists = []
        for item in self.by_type["AWS::CloudFront::Distribution"]:
            cfg = item.get("configuration", {})
            dist_cfg = cfg.get("distributionConfig", cfg)
            origins = dist_cfg.get("origins", {}).get("items", [])
            origin_domains = [o.get("domainName", "") for o in origins]
            dists.append({
                "id": item["resourceId"],
                "name": item.get("tags", {}).get("Name", dist_cfg.get("comment", item["resourceId"])),
                "domain_name": cfg.get("domainName", ""),
                "origin_domains": origin_domains,
                "waf_id": dist_cfg.get("webACLId", ""),
            })
        return dists

    def get_api_gateways(self):
        apis = []
        for item in self.by_type["AWS::ApiGateway::RestApi"]:
            cfg = item.get("configuration", {})
            apis.append({
                "id": item["resourceId"],
                "name": cfg.get("name", item.get("tags", {}).get("Name", item["resourceId"])),
                "type": "REST",
            })
        for item in self.by_type["AWS::ApiGatewayV2::Api"]:
            cfg = item.get("configuration", {})
            apis.append({
                "id": item["resourceId"],
                "name": cfg.get("name", item.get("tags", {}).get("Name", item["resourceId"])),
                "type": cfg.get("protocolType", "HTTP"),
            })
        return apis

    def get_route53_hosted_zones(self):
        zones = []
        for item in self.by_type["AWS::Route53::HostedZone"]:
            cfg = item.get("configuration", {})
            zones.append({
                "id": item["resourceId"],
                "name": cfg.get("name", item.get("tags", {}).get("Name", item["resourceId"])),
                "private": cfg.get("config", {}).get("privateZone", False),
            })
        return zones

    def get_dynamodb_tables(self):
        tables = []
        for item in self.by_type["AWS::DynamoDB::Table"]:
            cfg = item.get("configuration", {})
            tables.append({
                "id": item["resourceId"],
                "name": cfg.get("tableName", item.get("tags", {}).get("Name", item["resourceId"])),
                "status": cfg.get("tableStatus", ""),
            })
        return tables

    def get_elasticache_clusters(self):
        clusters = []
        for item in self.by_type["AWS::ElastiCache::CacheCluster"]:
            cfg = item.get("configuration", {})
            sg_ids = [sg.get("securityGroupId", "") for sg in cfg.get("securityGroups", [])]
            clusters.append({
                "id": item["resourceId"],
                "name": cfg.get("cacheClusterId", item.get("tags", {}).get("Name", item["resourceId"])),
                "engine": cfg.get("engine", ""),
                "node_type": cfg.get("cacheNodeType", ""),
                "sg_ids": sg_ids,
            })
        return clusters

    def get_redshift_clusters(self):
        clusters = []
        for item in self.by_type["AWS::Redshift::Cluster"]:
            cfg = item.get("configuration", {})
            sg_ids = [sg.get("vpcSecurityGroupId", "") for sg in cfg.get("vpcSecurityGroups", [])]
            clusters.append({
                "id": item["resourceId"],
                "name": cfg.get("clusterIdentifier", item.get("tags", {}).get("Name", item["resourceId"])),
                "node_type": cfg.get("nodeType", ""),
                "sg_ids": sg_ids,
            })
        return clusters

    def get_sqs_queues(self):
        queues = []
        for item in self.by_type["AWS::SQS::Queue"]:
            cfg = item.get("configuration", {})
            name = cfg.get("queueName", "")
            if not name:
                # Extract from ARN
                arn = item.get("arn", "")
                name = arn.split(":")[-1] if arn else item["resourceId"]
            queues.append({
                "id": item["resourceId"],
                "name": name,
                "fifo": cfg.get("fifoQueue", False),
            })
        return queues

    def get_sns_topics(self):
        topics = []
        for item in self.by_type["AWS::SNS::Topic"]:
            cfg = item.get("configuration", {})
            arn = cfg.get("topicArn", item.get("arn", ""))
            name = arn.split(":")[-1] if arn else item["resourceId"]
            topics.append({
                "id": item["resourceId"],
                "name": name,
            })
        return topics

    def get_kms_keys(self):
        keys = []
        for item in self.by_type["AWS::KMS::Key"]:
            cfg = item.get("configuration", {})
            keys.append({
                "id": item["resourceId"],
                "name": cfg.get("description", item.get("tags", {}).get("Name", item["resourceId"])),
                "state": cfg.get("keyState", ""),
            })
        return keys

    def get_cloudtrail_trails(self):
        trails = []
        for item in self.by_type["AWS::CloudTrail::Trail"]:
            cfg = item.get("configuration", {})
            trails.append({
                "id": item["resourceId"],
                "name": cfg.get("name", item.get("tags", {}).get("Name", item["resourceId"])),
                "s3_bucket": cfg.get("s3BucketName", ""),
                "is_logging": cfg.get("isLogging", False),
            })
        return trails

    def get_cloudwatch_alarms(self):
        alarms = []
        for item in self.by_type["AWS::CloudWatch::Alarm"]:
            cfg = item.get("configuration", {})
            alarms.append({
                "id": item["resourceId"],
                "name": cfg.get("alarmName", item.get("tags", {}).get("Name", item["resourceId"])),
                "metric": cfg.get("metricName", ""),
                "namespace": cfg.get("namespace", ""),
            })
        return alarms

    def get_elasticbeanstalk_environments(self):
        envs = []
        for item in self.by_type["AWS::ElasticBeanstalk::Environment"]:
            cfg = item.get("configuration", {})
            envs.append({
                "id": item["resourceId"],
                "name": cfg.get("environmentName", item.get("tags", {}).get("Name", item["resourceId"])),
                "app_name": cfg.get("applicationName", ""),
                "status": cfg.get("status", ""),
            })
        return envs

    def get_service_connections(self):
        """Infer connections for services that don't use Security Groups."""
        connections = []

        # CloudFront -> ALB/S3 (origin domains)
        alb_dns_map = {}
        for item in self.by_type["AWS::ElasticLoadBalancingV2::LoadBalancer"]:
            cfg = item.get("configuration", {})
            dns = cfg.get("dNSName", "")
            if dns:
                alb_dns_map[dns.lower()] = item["resourceId"]

        for dist in self.get_cloudfront_distributions():
            for origin in dist.get("origin_domains", []):
                origin_l = origin.lower()
                if "elb.amazonaws.com" in origin_l:
                    for dns, alb_id in alb_dns_map.items():
                        if dns in origin_l or origin_l in dns:
                            connections.append({
                                "from_type": "cloudfront", "from_id": dist["id"],
                                "to_type": "alb", "to_id": alb_id,
                                "label": "HTTPS",
                            })
                elif "s3" in origin_l:
                    connections.append({
                        "from_type": "cloudfront", "from_id": dist["id"],
                        "to_type": "s3", "to_id": origin,
                        "label": "",
                    })

        # CloudTrail -> S3
        for trail in self.get_cloudtrail_trails():
            if trail.get("s3_bucket"):
                connections.append({
                    "from_type": "cloudtrail", "from_id": trail["id"],
                    "to_type": "s3", "to_id": trail["s3_bucket"],
                    "label": "Logs",
                })

        return connections

    def build_sg_to_resources_map(self):
        """Build mapping: sg_id -> list of (resource_type, resource_id, resource_name)."""
        sg_map = defaultdict(list)

        # EC2 instances
        for item in self.by_type["AWS::EC2::Instance"]:
            cfg = item.get("configuration", {})
            tags = item.get("tags", {})
            for sg in cfg.get("securityGroups", []):
                sg_map[sg["groupId"]].append({
                    "type": "EC2",
                    "id": item["resourceId"],
                    "name": tags.get("Name", item["resourceId"]),
                    "prefix": "ec2_",
                })

        # ALBs
        for item in self.by_type["AWS::ElasticLoadBalancingV2::LoadBalancer"]:
            cfg = item.get("configuration", {})
            for sg_id in cfg.get("securityGroups", []):
                sg_map[sg_id].append({
                    "type": "ALB",
                    "id": item["resourceId"],
                    "name": cfg.get("loadBalancerName", ""),
                    "prefix": "alb_",
                })

        # RDS
        for item in self.by_type["AWS::RDS::DBInstance"]:
            cfg = item.get("configuration", {})
            for sg in cfg.get("vpcSecurityGroups", []):
                sg_map[sg["vpcSecurityGroupId"]].append({
                    "type": "RDS",
                    "id": item["resourceId"],
                    "name": cfg.get("dBInstanceIdentifier", ""),
                    "prefix": "rds_",
                })

        # Lambda (VPC-attached only)
        for item in self.by_type["AWS::Lambda::Function"]:
            cfg = item.get("configuration", {})
            vpc_cfg = cfg.get("vpcConfig", {})
            for sg_id in vpc_cfg.get("securityGroupIds", []):
                sg_map[sg_id].append({
                    "type": "Lambda",
                    "id": item["resourceId"],
                    "name": cfg.get("functionName", ""),
                    "prefix": "lambda_",
                })

        # ECS Services (awsvpc)
        for item in self.by_type["AWS::ECS::Service"]:
            cfg = item.get("configuration", {})
            net_cfg = cfg.get("networkConfiguration", {}).get("awsvpcConfiguration", {})
            for sg_id in net_cfg.get("securityGroups", []):
                sg_map[sg_id].append({
                    "type": "ECS",
                    "id": item["resourceId"],
                    "name": cfg.get("serviceName", ""),
                    "prefix": "ecs_",
                })

        # EKS
        for item in self.by_type["AWS::EKS::Cluster"]:
            cfg = item.get("configuration", {})
            vpc_cfg = cfg.get("resourcesVpcConfig", {})
            for sg_id in vpc_cfg.get("securityGroupIds", []):
                sg_map[sg_id].append({
                    "type": "EKS",
                    "id": item["resourceId"],
                    "name": cfg.get("name", ""),
                    "prefix": "eks_",
                })

        # ElastiCache
        for item in self.by_type["AWS::ElastiCache::CacheCluster"]:
            cfg = item.get("configuration", {})
            for sg in cfg.get("securityGroups", []):
                sg_map[sg.get("securityGroupId", "")].append({
                    "type": "ElastiCache",
                    "id": item["resourceId"],
                    "name": cfg.get("cacheClusterId", ""),
                    "prefix": "cache_",
                })

        # Redshift
        for item in self.by_type["AWS::Redshift::Cluster"]:
            cfg = item.get("configuration", {})
            for sg in cfg.get("vpcSecurityGroups", []):
                sg_map[sg.get("vpcSecurityGroupId", "")].append({
                    "type": "Redshift",
                    "id": item["resourceId"],
                    "name": cfg.get("clusterIdentifier", ""),
                    "prefix": "redshift_",
                })

        return sg_map


# ============================================================
# PPTX Diagram Generator
# ============================================================
class DiagramGenerator:
    """Generate PowerPoint network diagram from parsed AWS Config data."""

    def __init__(self, parser: AWSConfigParser):
        self.parser = parser
        self.prs = Presentation()
        # Widescreen 16:9
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

        # Layout tracking
        self.shape_positions = {}  # resource_id -> (center_x, center_y)

    def generate(self, output_path):
        """Main generation entry point."""
        vpcs = self.parser.get_vpcs()
        peerings = self.parser.get_peering_connections()

        # Slide 1: Overview diagram
        self._create_overview_slide(vpcs, peerings)

        # Slide 2: Security Group rules (audit detail)
        self._create_sg_detail_slide(vpcs)

        # Slide 3: Traffic flow summary
        self._create_traffic_flow_slide(vpcs)

        self.prs.save(output_path)
        print(f"Diagram saved: {output_path}")

    # ----------------------------------------------------------
    # Slide 1: Overview
    # ----------------------------------------------------------
    def _create_overview_slide(self, vpcs, peerings):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank

        # Title
        self._add_text_box(slide, Inches(0.3), Inches(0.15), Inches(10), Inches(0.5),
                           "AWS Network Architecture - External Audit Overview",
                           font_size=20, bold=True)

        # ---- Internet (top center) ----
        inet_x, inet_y = Inches(4.8), Inches(0.7)
        inet_w, inet_h = Inches(2.0), Inches(0.5)
        self._add_rounded_rect(slide, inet_x, inet_y, inet_w, inet_h,
                                "Internet", Colors.INTERNET_FILL, Colors.INTERNET_BORDER,
                                font_size=11, bold=True)
        self.shape_positions["internet"] = (int(inet_x + inet_w / 2), int(inet_y + inet_h / 2))

        # ---- WAF (between Internet and VPC) ----
        waf_drawn = False
        # Find ALBs to check for WAF association
        for vpc in vpcs:
            for alb in self.parser.get_albs_for_vpc(vpc["id"]):
                waf = self.parser.get_waf_for_alb(alb["id"])
                if waf:
                    waf_x, waf_y = Inches(4.5), Inches(1.35)
                    waf_w, waf_h = Inches(2.6), Inches(0.55)
                    rules_text = f"WAF: {waf['name']}\nRules: {waf['rule_count']}"
                    self._add_rounded_rect(slide, waf_x, waf_y, waf_w, waf_h,
                                            rules_text, Colors.WAF_FILL, Colors.WAF_BORDER,
                                            font_size=8, bold=True)
                    self.shape_positions["waf"] = (int(waf_x + waf_w / 2), int(waf_y + waf_h / 2))

                    # Internet -> WAF arrow
                    self._add_arrow(slide,
                                    inet_x + inet_w / 2, inet_y + inet_h,
                                    waf_x + waf_w / 2, waf_y,
                                    Colors.ARROW_EXTERNAL, width=Pt(2.5),
                                    label=":443 / :80")
                    waf_drawn = True
                    break
            if waf_drawn:
                break

        # ---- S3 (right side, outside VPC) ----
        s3_buckets = self.parser.get_s3_buckets()
        if s3_buckets:
            s3_x, s3_y = Inches(11.3), Inches(3.5)
            s3_w, s3_h = Inches(1.7), Inches(1.0)
            s3_texts = []
            for b in s3_buckets:
                enc = "KMS" if b["encrypted"] else "None"
                s3_texts.append(f"S3: {b['name']}\nEnc: {enc}\nVer: {b['versioning']}")
            s3_label = "\n".join(s3_texts)
            self._add_rounded_rect(slide, s3_x, s3_y, s3_w, s3_h,
                                    s3_label, Colors.S3_FILL, Colors.S3_BORDER,
                                    font_size=7, bold=False)
            self.shape_positions["s3"] = (int(s3_x + s3_w / 2), int(s3_y + s3_h / 2))

        # ---- VPCs ----
        vpc_start_y = Inches(2.1)
        vpc_gap = Inches(0.4)
        # Reserve right side for S3 if present
        total_vpc_width = Inches(10.7) if s3_buckets else Inches(12.5)

        if len(vpcs) == 1:
            vpc_width = total_vpc_width
            vpc_positions = [(Inches(0.4), vpc_start_y)]
        else:
            vpc_width = (total_vpc_width - vpc_gap * (len(vpcs) - 1)) / len(vpcs)
            vpc_positions = []
            for i in range(len(vpcs)):
                x = Inches(0.4) + i * (vpc_width + vpc_gap)
                vpc_positions.append((x, vpc_start_y))

        for i, vpc in enumerate(vpcs):
            vx, vy = vpc_positions[i]
            self._draw_vpc(slide, vpc, vx, vy, vpc_width, Inches(5.2))

        # ---- Peering connections ----
        for peer in peerings:
            self._draw_peering_arrow(slide, peer, vpcs, vpc_positions, vpc_width)

        # ---- WAF -> IGW or IGW -> Internet arrow ----
        for i, vpc in enumerate(vpcs):
            igw = self.parser.get_igw_for_vpc(vpc["id"])
            if igw:
                igw_key = f"igw_{igw['id']}"
                if igw_key in self.shape_positions:
                    sx, sy = self.shape_positions[igw_key]
                    if waf_drawn and "waf" in self.shape_positions:
                        # WAF -> IGW arrow
                        wx, wy = self.shape_positions["waf"]
                        self._add_arrow(slide, wx, wy + Inches(0.28),
                                        sx, sy - Inches(0.28),
                                        Colors.ARROW_EXTERNAL, width=Pt(2))
                    else:
                        # Direct IGW -> Internet
                        ex, ey = self.shape_positions["internet"]
                        self._add_arrow(slide, sx, sy - Inches(0.15),
                                        ex, ey + Inches(0.3),
                                        Colors.ARROW_EXTERNAL, width=Pt(2.5),
                                        label=":443 / :80")

        # ---- Internal traffic arrows (SG based) ----
        self._draw_traffic_arrows(slide)

        # ---- NAT -> Internet outbound arrow ----
        for vpc in vpcs:
            for nat in self.parser.get_nat_gateways_for_vpc(vpc["id"]):
                nat_key = f"nat_{nat['id']}"
                if nat_key in self.shape_positions:
                    nx, ny = self.shape_positions[nat_key]
                    ix, iy = self.shape_positions["internet"]
                    self._add_arrow(slide, nx, ny - Inches(0.28),
                                    ix + Inches(0.8), iy + Inches(0.25),
                                    RGBColor(0xFF, 0x8F, 0x00), width=Pt(1.5),
                                    label="Outbound :443")

        # ---- App -> S3 arrow ----
        if s3_buckets and "s3" in self.shape_positions:
            # Find any App server to draw arrow from
            for vpc in vpcs:
                subnets = self.parser.get_subnets_for_vpc(vpc["id"])
                for sub in subnets:
                    instances = self.parser.get_instances_for_subnet(sub["id"])
                    for inst in instances:
                        if inst.get("role") == "AppServer":
                            app_key = f"ec2_{inst['id']}"
                            if app_key in self.shape_positions:
                                ax, ay = self.shape_positions[app_key]
                                s3x, s3y = self.shape_positions["s3"]
                                self._add_arrow(slide,
                                                ax + Inches(0.75), ay,
                                                s3x - Inches(0.1), s3y,
                                                Colors.S3_BORDER, width=Pt(1.5),
                                                label="S3 API :443")
                                break  # one arrow is enough
                    else:
                        continue
                    break

        # ---- Legend ----
        legend_x = Inches(0.3) if s3_buckets else Inches(10.5)
        legend_y = Inches(0.15) if s3_buckets else Inches(0.15)
        self._add_legend(slide, Inches(10.5), Inches(0.15))

    def _draw_vpc(self, slide, vpc, x, y, w, h):
        """Draw a VPC with AZ-column layout. Traffic flows top→bottom, left→right.

        Layout strategy:
        - Left column: Infrastructure (IGW, NAT) — gateway services
        - AZ columns: AZ-a, AZ-c — per-AZ resources (EC2 instances)
        - Cross-AZ resources (ALB, RDS): centered spanning all AZ columns
        - Tiers ordered top→bottom: Public → Private → Isolated
        """
        # VPC box
        self._add_rounded_rect(
            slide, x, y, w, h,
            "", Colors.VPC_FILL, Colors.VPC_BORDER, font_size=10)

        # VPC header
        header = f"VPC: {vpc['name']}  |  {vpc['cidr']}  |  {vpc['region']}"
        self._add_text_box(slide, x + Inches(0.15), y + Inches(0.05),
                           w - Inches(0.3), Inches(0.35),
                           header, font_size=9, bold=True,
                           color=Colors.VPC_BORDER)

        # ---- Collect all data ----
        subnets = self.parser.get_subnets_for_vpc(vpc["id"])
        igw = self.parser.get_igw_for_vpc(vpc["id"])
        nats = self.parser.get_nat_gateways_for_vpc(vpc["id"])
        albs = self.parser.get_albs_for_vpc(vpc["id"])
        rds_list = self.parser.get_rds_for_vpc(vpc["id"])

        # Group subnets by tier
        tiers = defaultdict(list)
        for s in subnets:
            tiers[s["tier"]].append(s)

        # Discover AZ columns from subnets
        az_set = sorted(set(s["az"] for s in subnets if s["az"]))
        if not az_set:
            az_set = ["default"]
        num_az = len(az_set)
        az_index = {az: i for i, az in enumerate(az_set)}

        # ---- Layout dimensions ----
        margin = Inches(0.12)
        tier_order = ["Public", "Private", "Isolated"]
        active_tiers = [t for t in tier_order if t in tiers]
        tier_count = len(active_tiers)
        if tier_count == 0:
            return

        # Reserve left column for infra (IGW, NAT) — visible in all tiers
        infra_col_w = Inches(1.5)
        has_infra = bool(igw or nats)

        # Tier area (inner content area of VPC)
        tier_area_x = x + margin
        tier_area_y = y + Inches(0.42)
        tier_area_w = w - margin * 2
        tier_area_h = h - Inches(0.52)

        tier_gap = Inches(0.08)
        tier_h = (tier_area_h - tier_gap * (tier_count - 1)) / tier_count

        # AZ column area (right of infra column)
        if has_infra:
            az_area_x = tier_area_x + infra_col_w + Inches(0.08)
            az_area_w = tier_area_w - infra_col_w - Inches(0.08)
        else:
            az_area_x = tier_area_x + Inches(0.1)
            az_area_w = tier_area_w - Inches(0.2)

        az_col_gap = Inches(0.15)
        az_col_w = int((az_area_w - az_col_gap * (num_az - 1)) / num_az) if num_az > 0 else az_area_w

        res_w = min(Inches(1.9), int(az_col_w - Inches(0.15)))
        res_h = Inches(0.6)
        # Wider resource box for cross-AZ (ALB, RDS)
        cross_az_w = min(Inches(2.4), int(az_area_w - Inches(0.2)))

        # ---- AZ column headers ----
        for i, az in enumerate(az_set):
            az_short = az.split("-")[-1] if "-" in az else az  # e.g. "1a"
            hdr_x = int(az_area_x + i * (az_col_w + az_col_gap))
            self._add_text_box(slide, hdr_x, tier_area_y - Inches(0.02),
                               az_col_w, Inches(0.2),
                               f"AZ: {az_short}", font_size=7, bold=True,
                               color=Colors.TEXT_GRAY, align=PP_ALIGN.CENTER)

        # ---- Draw each tier ----
        current_y = tier_area_y + Inches(0.15)

        for tier in active_tiers:
            tier_subnets = tiers[tier]
            fill, border = self._tier_colors(tier)

            # Tier container
            self._add_rounded_rect(
                slide, tier_area_x, current_y,
                tier_area_w, tier_h,
                "", fill, border, font_size=8)

            # Tier label
            tier_label = f"{tier} Subnet"
            if tier == "Public":
                tier_label += "  [EXTERNAL FACING]"
            self._add_text_box(
                slide, tier_area_x + Inches(0.08), current_y + Inches(0.02),
                tier_area_w - Inches(0.2), Inches(0.2),
                tier_label, font_size=7, bold=True, color=border)

            res_start_y = current_y + Inches(0.28)

            # ---- Infrastructure column (left, only in Public tier) ----
            if tier == "Public" and has_infra:
                infra_y = res_start_y
                infra_x = int(tier_area_x + Inches(0.08))
                infra_item_w = int(infra_col_w - Inches(0.16))

                if igw:
                    self._add_rounded_rect(
                        slide, infra_x, infra_y, infra_item_w, Inches(0.45),
                        f"IGW\n{igw['name']}", Colors.IGW_FILL, Colors.IGW_BORDER,
                        font_size=7, bold=True)
                    self.shape_positions[f"igw_{igw['id']}"] = (
                        int(infra_x + infra_item_w / 2), int(infra_y + Inches(0.225)))
                    infra_y += Inches(0.52)

                for nat in nats:
                    label = f"NAT GW\n{nat['public_ip']}"
                    self._add_rounded_rect(
                        slide, infra_x, infra_y, infra_item_w, Inches(0.45),
                        label, Colors.NAT_FILL, Colors.NAT_BORDER,
                        font_size=7, bold=True)
                    self.shape_positions[f"nat_{nat['id']}"] = (
                        int(infra_x + infra_item_w / 2), int(infra_y + Inches(0.225)))
                    infra_y += Inches(0.52)

            # ---- Cross-AZ resources (ALB, RDS) - centered in AZ area ----
            cross_az_items = []  # (label, rid, fill, border, prefix)

            if tier == "Public":
                for alb in albs:
                    if alb["scheme"] == "internet-facing":
                        label = f"ALB: {alb['name']}  (internet-facing, spans AZs)"
                        cross_az_items.append((label, alb["id"],
                                               Colors.ALB_FILL, Colors.ALB_BORDER, "alb_"))

            if tier in ("Isolated", "Private"):
                for db in rds_list:
                    tier_subnet_ids = {s["id"] for s in tier_subnets}
                    if tier_subnet_ids & set(db["subnet_ids"]):
                        multi = " (Multi-AZ)" if db["multi_az"] else ""
                        label = f"RDS {db['engine']}{multi}\n{db['name']}  |  Port: {db['port']}"
                        cross_az_items.append((label, db["id"],
                                               Colors.RDS_FILL, Colors.RDS_BORDER, "rds_"))

            # Draw cross-AZ items centered
            cross_y = res_start_y
            cross_center_x = int(az_area_x + az_area_w / 2 - cross_az_w / 2)
            for label, rid, cfill, cborder, prefix in cross_az_items:
                self._add_rounded_rect(
                    slide, cross_center_x, cross_y, cross_az_w, res_h,
                    label, cfill, cborder, font_size=7, bold=True)
                self.shape_positions[f"{prefix}{rid}"] = (
                    int(cross_center_x + cross_az_w / 2), int(cross_y + res_h / 2))
                cross_y += res_h + Inches(0.06)

            # ---- AZ-aligned resources (per-AZ EC2 instances) ----
            # Start after cross-AZ items
            per_az_start_y = cross_y if cross_az_items else res_start_y

            for col_idx in range(num_az):
                col_x = int(az_area_x + col_idx * (az_col_w + az_col_gap))
                res_center_x = int(col_x + az_col_w / 2 - res_w / 2)
                ry = per_az_start_y

                # EC2 instances for this AZ in this tier
                for sub in tier_subnets:
                    if az_index.get(sub["az"], 0) != col_idx:
                        continue
                    instances = self.parser.get_instances_for_subnet(sub["id"])
                    for inst in instances:
                        role = f" ({inst['role']})" if inst['role'] else ""
                        label = f"EC2{role}\n{inst['name']}\n{inst['private_ip']}"
                        self._add_rounded_rect(
                            slide, res_center_x, ry, res_w, res_h,
                            label, Colors.EC2_FILL, Colors.EC2_BORDER, font_size=7)
                        self.shape_positions[f"ec2_{inst['id']}"] = (
                            int(res_center_x + res_w / 2), int(ry + res_h / 2))
                        ry += res_h + Inches(0.06)

            current_y += tier_h + tier_gap

    def _draw_peering_arrow(self, slide, peer, vpcs, vpc_positions, vpc_width):
        """Draw VPC Peering connection arrow."""
        # Find positions of accepter and requester VPCs
        acc_idx = req_idx = None
        for i, vpc in enumerate(vpcs):
            if vpc["id"] == peer["accepter_vpc"]:
                acc_idx = i
            if vpc["id"] == peer["requester_vpc"]:
                req_idx = i

        if acc_idx is not None and req_idx is not None:
            ax, ay = vpc_positions[acc_idx]
            rx, ry = vpc_positions[req_idx]

            # Arrow from right edge of left VPC to left edge of right VPC
            if acc_idx < req_idx:
                sx = ax + vpc_width
                ex = rx
            else:
                sx = rx + vpc_width
                ex = ax

            mid_y = ay + Inches(2.5)
            self._add_arrow(slide, sx, mid_y, ex, mid_y,
                            Colors.ARROW_PEERING, width=Pt(2),
                            label=f"Peering: {peer['name']}")

    def _draw_traffic_arrows(self, slide):
        """Draw arrows between resources based on SG ingress/egress rules.

        Strategy to minimize crossing:
        - For 1-to-many or many-to-many: prefer same-column (vertical) arrows
        - Cross-column (diagonal) arrows are avoided when both sides have
          resources in the same column. Only draw cross-column if no same-column
          pair exists.
        - Arrows flow downward (top→bottom) following tier order.
        """
        sg_connections = self.parser.get_sg_connections()
        sg_resource_map = self.parser.build_sg_to_resources_map()

        drawn = set()

        for conn in sg_connections:
            from_sg = conn["from_sg"]
            to_sg = conn["to_sg"]
            port = conn.get("port", "")

            from_resources = sg_resource_map.get(from_sg, [])
            to_resources = sg_resource_map.get(to_sg, [])

            # Build all possible arrows with position info
            arrows = []
            for from_res in from_resources:
                from_key = f"{from_res['prefix']}{from_res['id']}"
                if from_key not in self.shape_positions:
                    continue
                sx, sy = self.shape_positions[from_key]

                for to_res in to_resources:
                    to_key = f"{to_res['prefix']}{to_res['id']}"
                    if to_key not in self.shape_positions:
                        continue
                    ex, ey = self.shape_positions[to_key]

                    # Classify: same-column (vertical) vs cross-column (diagonal)
                    x_diff = abs(sx - ex)
                    is_same_col = x_diff < Inches(1.0)

                    arrows.append({
                        "from_key": from_key, "to_key": to_key,
                        "sx": sx, "sy": sy, "ex": ex, "ey": ey,
                        "same_col": is_same_col,
                    })

            # If we have same-column arrows, skip cross-column ones
            has_same_col = any(a["same_col"] for a in arrows)

            for arrow in arrows:
                # Skip cross-column if same-column arrows exist
                if has_same_col and not arrow["same_col"]:
                    continue

                from_key = arrow["from_key"]
                to_key = arrow["to_key"]
                arrow_id = f"{from_key}->{to_key}"
                if arrow_id in drawn:
                    continue
                drawn.add(arrow_id)

                sx, sy = arrow["sx"], arrow["sy"]
                ex, ey = arrow["ex"], arrow["ey"]

                color = Colors.ARROW_INTERNAL
                label = f":{port}" if port else ""

                # Adjust start/end to shape edges
                if abs(ey - sy) > abs(ex - sx):
                    # Vertical arrow (most common: tier-to-tier)
                    if ey > sy:
                        sy_adj = sy + Inches(0.30)
                        ey_adj = ey - Inches(0.30)
                    else:
                        sy_adj = sy - Inches(0.30)
                        ey_adj = ey + Inches(0.30)
                    self._add_arrow(slide, sx, int(sy_adj), ex, int(ey_adj),
                                    color, width=Pt(1.8), label=label)
                else:
                    # Horizontal or diagonal
                    if ex > sx:
                        sx_adj = sx + Inches(0.85)
                        ex_adj = ex - Inches(0.85)
                    else:
                        sx_adj = sx - Inches(0.85)
                        ex_adj = ex + Inches(0.85)
                    self._add_arrow(slide, int(sx_adj), sy, int(ex_adj), ey,
                                    color, width=Pt(1.8), label=label)

    # ----------------------------------------------------------
    # Slide 2: Security Group Details
    # ----------------------------------------------------------
    def _create_sg_detail_slide(self, vpcs):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        self._add_text_box(slide, Inches(0.3), Inches(0.2), Inches(12), Inches(0.5),
                           "Security Group Rules - External Access Audit",
                           font_size=20, bold=True)

        y = Inches(0.9)

        for vpc in vpcs:
            # VPC header
            self._add_text_box(slide, Inches(0.3), y, Inches(6), Inches(0.35),
                               f"VPC: {vpc['name']} ({vpc['cidr']})",
                               font_size=12, bold=True, color=Colors.VPC_BORDER)
            y += Inches(0.4)

            # External rules (0.0.0.0/0) - highlighted in red
            ext_rules = self.parser.get_external_sg_rules(vpc["id"])
            if ext_rules:
                self._add_text_box(slide, Inches(0.5), y, Inches(6), Inches(0.25),
                                   "!! EXTERNAL ACCESS (0.0.0.0/0) !!",
                                   font_size=10, bold=True, color=Colors.ALERT_RED)
                y += Inches(0.3)

                # Table header
                headers = ["SG Name", "Direction", "Port", "Source", "Description"]
                col_widths = [Inches(2), Inches(1.2), Inches(1), Inches(1.5), Inches(4)]
                self._draw_table_row(slide, Inches(0.5), y, col_widths, headers,
                                     bold=True, bg_color=RGBColor(0xFF, 0xCC, 0xCC))
                y += Inches(0.3)

                for rule in ext_rules:
                    port_str = str(rule["from_port"])
                    if rule["from_port"] != rule["to_port"]:
                        port_str = f"{rule['from_port']}-{rule['to_port']}"
                    row = [
                        rule["sg_name"],
                        rule["direction"],
                        port_str,
                        rule["source"],
                        rule["description"],
                    ]
                    self._draw_table_row(slide, Inches(0.5), y, col_widths, row,
                                         bg_color=RGBColor(0xFF, 0xF0, 0xF0))
                    y += Inches(0.28)
            else:
                self._add_text_box(slide, Inches(0.5), y, Inches(6), Inches(0.25),
                                   "No external access rules (0.0.0.0/0)",
                                   font_size=9, color=Colors.TEXT_GRAY)
                y += Inches(0.3)

            # All SG rules summary
            y += Inches(0.1)
            self._add_text_box(slide, Inches(0.5), y, Inches(6), Inches(0.25),
                               "All Security Group Rules:",
                               font_size=10, bold=True)
            y += Inches(0.3)

            sgs = self.parser.get_security_groups_for_vpc(vpc["id"])
            for sg in sgs:
                self._add_text_box(slide, Inches(0.5), y, Inches(8), Inches(0.22),
                                   f"{sg['name']} ({sg['id']}) - {sg['description']}",
                                   font_size=8, bold=True)
                y += Inches(0.25)

                # Ingress
                for rule in sg.get("ingress", []):
                    sources = []
                    for ipr in AWSConfigParser._normalize_ip_ranges(rule.get("ipRanges", [])):
                        sources.append(ipr.get("cidrIp", ""))
                    for pair in rule.get("userIdGroupPairs", []):
                        sources.append(f"SG:{pair.get('groupId', '')}")

                    port_str = str(rule.get("fromPort", "all"))
                    src_str = ", ".join(sources) if sources else "N/A"
                    is_external = any("0.0.0.0/0" in s for s in sources)

                    text = f"  IN  | :{port_str} | from {src_str}"
                    color = Colors.ALERT_RED if is_external else Colors.TEXT_BLACK
                    self._add_text_box(slide, Inches(0.7), y, Inches(10), Inches(0.2),
                                       text, font_size=7, color=color,
                                       bold=is_external)
                    y += Inches(0.2)

                # Egress
                for rule in sg.get("egress", []):
                    dests = []
                    for ipr in AWSConfigParser._normalize_ip_ranges(rule.get("ipRanges", [])):
                        dests.append(ipr.get("cidrIp", ""))
                    for pair in rule.get("userIdGroupPairs", []):
                        dests.append(f"SG:{pair.get('groupId', '')}")

                    port_str = str(rule.get("fromPort", "all"))
                    dst_str = ", ".join(dests) if dests else "N/A"

                    text = f"  OUT | :{port_str} | to {dst_str}"
                    self._add_text_box(slide, Inches(0.7), y, Inches(10), Inches(0.2),
                                       text, font_size=7)
                    y += Inches(0.2)

                y += Inches(0.1)

            y += Inches(0.2)

    # ----------------------------------------------------------
    # Slide 3: Traffic Flow
    # ----------------------------------------------------------
    def _create_traffic_flow_slide(self, vpcs):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        self._add_text_box(slide, Inches(0.3), Inches(0.2), Inches(12), Inches(0.5),
                           "Traffic Flow Summary - External Audit",
                           font_size=20, bold=True)

        # Get SG-to-SG connections
        connections = self.parser.get_sg_connections()

        # Build a readable flow
        y = Inches(1.0)

        # Inbound flows from Internet
        self._add_text_box(slide, Inches(0.3), y, Inches(8), Inches(0.35),
                           "INBOUND Traffic Flows (Internet -> Internal)",
                           font_size=14, bold=True, color=Colors.ARROW_EXTERNAL)
        y += Inches(0.5)

        for vpc in vpcs:
            ext_rules = self.parser.get_external_sg_rules(vpc["id"])
            if ext_rules:
                self._add_text_box(slide, Inches(0.5), y, Inches(8), Inches(0.25),
                                   f"VPC: {vpc['name']}", font_size=11, bold=True)
                y += Inches(0.35)

                # Trace the chain: Internet -> ALB (SG with 0.0.0.0/0) -> Web -> DB
                for rule in ext_rules:
                    text = f"Internet -> :{rule['from_port']} -> {rule['sg_name']}"
                    self._add_text_box(slide, Inches(0.7), y, Inches(10), Inches(0.22),
                                       text, font_size=9, color=Colors.ALERT_RED)
                    y += Inches(0.25)

                    # Follow egress from this SG
                    for conn in connections:
                        if conn["from_sg"] == rule["sg_id"]:
                            text2 = f"    -> :{conn['port']} -> {conn['to_sg_name']} ({conn['description']})"
                            self._add_text_box(slide, Inches(0.7), y, Inches(10), Inches(0.22),
                                               text2, font_size=9)
                            y += Inches(0.25)

                            # Follow next hop
                            for conn2 in connections:
                                if conn2["from_sg"] == conn["to_sg"]:
                                    text3 = f"        -> :{conn2['port']} -> {conn2['to_sg_name']} ({conn2['description']})"
                                    self._add_text_box(slide, Inches(0.7), y, Inches(10), Inches(0.22),
                                                       text3, font_size=9)
                                    y += Inches(0.25)

        y += Inches(0.3)

        # VPC Peering
        peerings = self.parser.get_peering_connections()
        if peerings:
            self._add_text_box(slide, Inches(0.3), y, Inches(8), Inches(0.35),
                               "VPC Peering Connections",
                               font_size=14, bold=True, color=Colors.ARROW_PEERING)
            y += Inches(0.5)

            for peer in peerings:
                text = (f"{peer['name']}: "
                        f"{peer['requester_vpc']} ({peer['requester_cidr']}) "
                        f"<-> {peer['accepter_vpc']} ({peer['accepter_cidr']})")
                self._add_text_box(slide, Inches(0.5), y, Inches(10), Inches(0.22),
                                   text, font_size=9)
                y += Inches(0.3)

        # Outbound flows
        y += Inches(0.2)
        self._add_text_box(slide, Inches(0.3), y, Inches(8), Inches(0.35),
                           "OUTBOUND Traffic (Internal -> Internet)",
                           font_size=14, bold=True, color=RGBColor(0xFF, 0x8F, 0x00))
        y += Inches(0.5)

        for vpc in vpcs:
            sgs = self.parser.get_security_groups_for_vpc(vpc["id"])
            has_outbound = False
            for sg in sgs:
                for rule in sg.get("egress", []):
                    for ipr in AWSConfigParser._normalize_ip_ranges(rule.get("ipRanges", [])):
                        if ipr.get("cidrIp") == "0.0.0.0/0":
                            if not has_outbound:
                                self._add_text_box(slide, Inches(0.5), y, Inches(8), Inches(0.25),
                                                   f"VPC: {vpc['name']}", font_size=11, bold=True)
                                y += Inches(0.35)
                                has_outbound = True
                            text = (f"{sg['name']} -> :{rule.get('fromPort', 'all')} -> "
                                    f"0.0.0.0/0 ({ipr.get('description', '')})")
                            self._add_text_box(slide, Inches(0.7), y, Inches(10), Inches(0.22),
                                               text, font_size=9)
                            y += Inches(0.25)

    # ----------------------------------------------------------
    # Drawing Helpers
    # ----------------------------------------------------------
    def _add_rounded_rect(self, slide, x, y, w, h, text, fill_color, border_color,
                           font_size=10, bold=False):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)

        # Adjust corner radius
        shape.adjustments[0] = 0.05

        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        tf.margin_left = Pt(4)
        tf.margin_right = Pt(4)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)

        if text:
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = Colors.TEXT_BLACK
            p.alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].space_before = Pt(0)
            tf.paragraphs[0].space_after = Pt(0)

        return shape

    def _add_text_box(self, slide, x, y, w, h, text, font_size=10,
                       bold=False, color=None, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color or Colors.TEXT_BLACK
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        return txBox

    def _add_arrow(self, slide, x1, y1, x2, y2, color, width=Pt(1.5), label=None):
        """Add an arrow connector with arrowhead via direct XML manipulation."""
        from pptx.oxml.ns import qn
        from lxml import etree

        # Ensure coordinates are integers (EMU units)
        x1, y1, x2, y2 = int(x1), int(y1), int(x2), int(y2)

        connector = slide.shapes.add_connector(
            1,  # straight connector
            x1, y1, x2, y2
        )
        connector.line.color.rgb = color
        connector.line.width = width

        # Add arrowhead via XML (most reliable method for python-pptx)
        cxnSp = connector._element
        spPr = cxnSp.find(qn('p:spPr'))
        if spPr is None:
            spPr = cxnSp.find(qn('p:cxnSp/p:spPr'))

        ln = spPr.find(qn('a:ln'))
        if ln is None:
            ln = etree.SubElement(spPr, qn('a:ln'))

        # End arrow (triangle)
        tailEnd = ln.find(qn('a:tailEnd'))
        if tailEnd is None:
            tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
        tailEnd.set('type', 'triangle')
        tailEnd.set('w', 'med')
        tailEnd.set('len', 'med')

        # Add label near midpoint
        if label:
            mid_x = (x1 + x2) / 2
            mid_y = (y1 + y2) / 2

            # Offset label slightly to not overlap the line
            offset_x = Inches(0.05)
            offset_y = -Inches(0.18)

            # For mostly vertical lines, offset to the right
            if abs(y2 - y1) > abs(x2 - x1):
                offset_x = Inches(0.15)
                offset_y = -Inches(0.05)

            self._add_text_box(slide,
                               mid_x + offset_x - Inches(0.35),
                               mid_y + offset_y,
                               Inches(0.9), Inches(0.22),
                               label, font_size=7, color=color,
                               align=PP_ALIGN.CENTER, bold=True)

    def _draw_table_row(self, slide, x, y, col_widths, values,
                         bold=False, bg_color=None):
        """Draw a table-like row using rectangles."""
        cx = x
        for i, (val, cw) in enumerate(zip(values, col_widths)):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, y, cw, Inches(0.28))
            if bg_color:
                shape.fill.solid()
                shape.fill.fore_color.rgb = bg_color
            else:
                shape.fill.background()
            shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            shape.line.width = Pt(0.5)

            tf = shape.text_frame
            tf.margin_left = Pt(3)
            tf.margin_right = Pt(3)
            tf.margin_top = Pt(1)
            tf.margin_bottom = Pt(1)
            p = tf.paragraphs[0]
            p.text = str(val)
            p.font.size = Pt(7)
            p.font.bold = bold
            p.alignment = PP_ALIGN.LEFT

            cx += cw

    def _add_legend(self, slide, x, y):
        """Add color legend."""
        items = [
            ("Public Subnet", Colors.PUBLIC_FILL, Colors.PUBLIC_BORDER),
            ("Private Subnet", Colors.PRIVATE_FILL, Colors.PRIVATE_BORDER),
            ("Isolated Subnet", Colors.ISOLATED_FILL, Colors.ISOLATED_BORDER),
            ("WAF", Colors.WAF_FILL, Colors.WAF_BORDER),
            ("Internet GW", Colors.IGW_FILL, Colors.IGW_BORDER),
            ("NAT Gateway", Colors.NAT_FILL, Colors.NAT_BORDER),
            ("ALB", Colors.ALB_FILL, Colors.ALB_BORDER),
            ("EC2", Colors.EC2_FILL, Colors.EC2_BORDER),
            ("RDS", Colors.RDS_FILL, Colors.RDS_BORDER),
            ("S3", Colors.S3_FILL, Colors.S3_BORDER),
        ]

        self._add_text_box(slide, x, y, Inches(2), Inches(0.25),
                           "Legend:", font_size=8, bold=True)
        ly = y + Inches(0.25)

        for label, fill, border in items:
            self._add_rounded_rect(slide, x, ly, Inches(0.3), Inches(0.2),
                                    "", fill, border, font_size=6)
            self._add_text_box(slide, x + Inches(0.35), ly, Inches(1.5), Inches(0.2),
                               label, font_size=7)
            ly += Inches(0.22)

        # Arrow legends
        ly += Inches(0.05)
        arrow_items = [
            ("External Access", Colors.ARROW_EXTERNAL),
            ("Internal Traffic", Colors.ARROW_INTERNAL),
            ("NAT Outbound", RGBColor(0xFF, 0x8F, 0x00)),
            ("VPC Peering", Colors.ARROW_PEERING),
            ("S3 Access", Colors.S3_BORDER),
        ]
        for label, color in arrow_items:
            self._add_text_box(slide, x, ly, Inches(0.3), Inches(0.2),
                               "-->", font_size=7, color=color, bold=True)
            self._add_text_box(slide, x + Inches(0.35), ly, Inches(1.5), Inches(0.2),
                               label, font_size=7)
            ly += Inches(0.22)

    def _tier_colors(self, tier):
        if tier == "Public":
            return Colors.PUBLIC_FILL, Colors.PUBLIC_BORDER
        elif tier == "Isolated":
            return Colors.ISOLATED_FILL, Colors.ISOLATED_BORDER
        else:
            return Colors.PRIVATE_FILL, Colors.PRIVATE_BORDER


# ============================================================
# Main
# ============================================================
def main():
    if len(sys.argv) < 2:
        print("Usage: python aws_config_parser.py <aws_config_snapshot.json>")
        print("Output: network_diagram.pptx (same directory as input)")
        sys.exit(1)

    input_path = sys.argv[1]
    if not os.path.exists(input_path):
        print(f"Error: File not found: {input_path}")
        sys.exit(1)

    output_dir = os.path.dirname(os.path.abspath(input_path))
    output_path = os.path.join(output_dir, "network_diagram.pptx")

    print(f"Parsing: {input_path}")
    parser = AWSConfigParser(input_path)

    print(f"Found resources:")
    for rtype, items in sorted(parser.by_type.items()):
        print(f"  {rtype}: {len(items)}")

    print(f"\nGenerating diagram...")
    generator = DiagramGenerator(parser)
    generator.generate(output_path)
    print("Done!")


if __name__ == "__main__":
    main()
